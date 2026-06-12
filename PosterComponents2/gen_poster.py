"""
gen_poster.py
─────────────
Builds the DAMO 699 capstone poster in TWO formats from one layout:

  • capstone_poster.svg   — vector mockup, open in any browser to preview
  • capstone_poster.pptx  — 48 x 36 in, fully editable in PowerPoint,
                            then "Save as PDF" for the printer (TPH wants PDF)

Conforms to the poster guidelines (9 required sections), kept high-level and
summarized. Edit the CONTENT block below to change any wording or numbers.

Run:  python gen_poster.py
"""

from __future__ import annotations
import base64, mimetypes
from pathlib import Path
from html import escape

HERE = Path(__file__).resolve().parent
ASSETS = HERE.parent / "PosterSamples"
LOGO   = ASSETS / "unfc_logo.png"
MASCOT = ASSETS / "unfc_mascot.jpeg"

def data_uri(p: Path) -> str:
    mime = mimetypes.guess_type(str(p))[0] or "image/png"
    return f"data:{mime};base64,{base64.b64encode(p.read_bytes()).decode()}"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = "14233F"; BLUE   = "2C6EBD"; TEAL   = "0E9384"; TEALB = "14B8A6"
AMBER  = "F59E0B"; AMBERD = "B5740A"; INK    = "1B2A4A"; MUTE  = "5A6B85"
BORDER = "D6E0EE"; WHITE  = "FFFFFF"; LIGHT  = "F4F7FB"
LBLUE  = "EAF1FB"; LTEAL  = "EAF6F4"; LAMBER = "FEF4E6"; LRED = "FDE7DE"
GREY   = "9AA8BE"

def hx(s): return "#" + s

# ── Poster geometry (inches) ─────────────────────────────────────────────────
PW, PH = 48.0, 36.0
COLW, GAP = 15.0, 0.65
C1, C2, C3 = 0.8, 16.45, 32.10
HH = 1.10                      # card header height

# ── CONTENT ──────────────────────────────────────────────────────────────────
TITLE    = ("Predict, Plan Ahead, Dispatch:",
            "A Unified Framework for Multi-Tenant Cluster Management with Overcommit and Guaranteed SLA")
AUTHOR   = "Tha Pyay Hmu  •  Lhagii Tsogtbayar  •  Nadia Rios  •  Jorge Mendoza  •  Alrick Grandison"
COURSE   = "DAMO 699: Capstone Project"
INSTR    = "Supervisor: PhD. Hany Osman"
INSTN    = "Master of Data Analytics, University of Niagara Falls"
DATE     = "June 2026"

SEC_PROBLEM = ("Research Problem & Objective", BLUE, [
    "Cloud clusters are shared by many tenants; reserving peak capacity for each "
    "tenant wastes memory and money.",
    "Objective: an optimization layer that safely overcommits shared machines while "
    "guaranteeing service level agreements (SLAs).",
    "It turns predicted demand into real-time and plan-ahead scheduling decisions.",
])
SEC_RQ = ("Research Questions & Hypotheses", NAVY, None)   # custom render
SEC_DATA = ("Data & Preprocessing", TEAL, [
    "Source: Google Cluster Usage Traces v3 (Borg), a subset of 9 tenants.",
    "All resources normalized to fractions of the largest machine (0 to 1 scale).",
    "Derived per-tenant statistics: job memory, CPU, arrival rate, and lifetime.",
    "Cleaned outliers, capped out-of-range values, and built 6-hour demand profiles.",
])
SEC_APPROACH = ("Analytical Approach", BLUE, [
    "Two-layer system: a prediction layer forecasts demand; the optimization layer "
    "decides placement.",
    "Real-time scheduler: mixed-integer optimization places jobs each interval.",
    "Plan-ahead allocator: chance-constrained optimization reserves capacity over a "
    "24-hour horizon.",
    "Justification: optimization guarantees feasibility and SLAs, unlike greedy rules.",
    "Diagnostics: solve time, placement rate, and SLA violation rate by problem size.",
])
SEC_IMPL = ("Practical Implications", TEAL, [
    "Operators serve more tenants per machine, cutting hardware cost.",
    "The SLA guarantee protects customer trust even while overcommitting.",
    "Plan-ahead reservations support capacity planning and budgeting.",
])
SEC_LIMITS = ("Limitations & Future Work", AMBERD, [
    "Based on a 9-tenant subset; production cells are larger and more diverse.",
    "Parameters are static and usage is assumed close to normal for the safety bound.",
    "Future: live retraining, more solver backends, broader tenant mix, online "
    "adaptation.",
])

KPIS = [   # (value, label, color)
    ("98%",  "Job placement rate",          BLUE),
    ("<10s", "Real-time solve / interval",  TEAL),
    ("≤10%", "SLA violation bound (ε)", AMBERD),
    ("2.4×", "More memory packed",     BLUE),
    ("9",    "Tenants co-scheduled",        TEAL),
]

RQ_TEXT = ("Can an optimization layer pack many tenants onto a shared cluster with "
           "overcommit while still guaranteeing SLAs at scale?")
H0_TEXT = ("Overcommit scheduling does not improve utilization without breaching SLAs "
           "versus a no-overcommit baseline.")
H1_TEXT = ("A chance-constrained optimization layer raises utilization while bounding "
           "SLA violations to a small target ε, and holds up as the cluster grows.")

# Chart data
BAR_CATS = ["No overcommit", "With framework"]
BAR_VALS = [55, 88]
LINE_HRS = ["0", "4", "8", "12", "16", "20", "24"]
LINE_USE = [40, 58, 74, 86, 80, 64, 50]
LINE_CAP = [90, 90, 90, 90, 90, 90, 90]


# ═════════════════════════════════════════════════════════════════════════════
#  SVG RENDERER  (100 px per inch; open in a browser)
# ═════════════════════════════════════════════════════════════════════════════
def build_svg() -> str:
    U = 100.0
    def px(v): return round(v * U, 1)
    out = []
    out.append(f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {px(PW)} {px(PH)}" '
               f'font-family="Arial, sans-serif">')
    out.append(f'<rect x="0" y="0" width="{px(PW)}" height="{px(PH)}" fill="#EEF2F8"/>')

    # ---- Banner ----
    bx, by, bw, bh = 0.8, 0.6, 46.4, 4.6
    out.append(f'<defs><linearGradient id="bg" x1="0" y1="0" x2="1" y2="0">'
               f'<stop offset="0" stop-color="{hx(NAVY)}"/>'
               f'<stop offset="0.7" stop-color="#1C3461"/>'
               f'<stop offset="1" stop-color="{hx(BLUE)}"/></linearGradient></defs>')
    out.append(f'<rect x="{px(bx)}" y="{px(by)}" width="{px(bw)}" height="{px(bh)}" rx="16" fill="url(#bg)"/>')
    out.append(f'<rect x="{px(bx)}" y="{px(by+bh-0.18)}" width="{px(bw/2)}" height="18" fill="{hx(AMBER)}"/>')
    out.append(f'<rect x="{px(bx+bw/2)}" y="{px(by+bh-0.18)}" width="{px(bw/2)}" height="18" fill="{hx(TEALB)}"/>')
    # real UNF logo (left) on a white rounded chip + fox mascot (right)
    out.append(f'<rect x="{px(1.2)}" y="{px(1.2)}" width="{px(2.8)}" height="{px(2.8)}" rx="14" fill="#FFFFFF"/>')
    out.append(f'<image x="{px(1.45)}" y="{px(1.45)}" width="{px(2.3)}" height="{px(2.3)}" '
               f'href="{data_uri(LOGO)}" preserveAspectRatio="xMidYMid meet"/>')
    out.append(f'<image x="{px(43.6)}" y="{px(0.85)}" width="{px(3.1)}" height="{px(3.1)}" '
               f'href="{data_uri(MASCOT)}" preserveAspectRatio="xMidYMid meet" '
               f'style="clip-path: inset(0 round 16px)"/>')
    # title text
    cx = bx + bw / 2
    out.append(f'<text x="{px(cx)}" y="{px(1.8)}" fill="#FFFFFF" font-size="64" font-weight="700" text-anchor="middle">{escape(TITLE[0])}</text>')
    out.append(f'<text x="{px(cx)}" y="{px(2.7)}" fill="#9FE7DA" font-size="40" font-weight="600" text-anchor="middle">{escape(TITLE[1])}</text>')
    out.append(f'<text x="{px(cx)}" y="{px(3.55)}" fill="#CBD8EE" font-size="30" text-anchor="middle">{escape(AUTHOR)}</text>')
    out.append(f'<text x="{px(cx)}" y="{px(4.25)}" fill="#9FB4D6" font-size="26" text-anchor="middle">'
               f'{escape(COURSE + "   |   " + INSTR + "   |   " + INSTN + "   |   " + DATE)}</text>')

    # ---- Card helper ----
    def card(x, y, w, h, num, title, color):
        X, Y, W, H = px(x), px(y), px(w), px(h)
        r = 18; hh = px(HH)
        out.append(f'<rect x="{X}" y="{Y}" width="{W}" height="{H}" rx="{r}" fill="#FFFFFF" stroke="{hx(BORDER)}" stroke-width="2"/>')
        out.append(f'<path d="M{X} {Y+r} a{r} {r} 0 0 1 {r} -{r} h{W-2*r} a{r} {r} 0 0 1 {r} {r} v{hh-r} h-{W} z" fill="{hx(color)}"/>')
        out.append(f'<circle cx="{X+58}" cy="{Y+55}" r="34" fill="#FFFFFF"/>')
        out.append(f'<text x="{X+58}" y="{Y+68}" text-anchor="middle" font-size="36" font-weight="800" fill="{hx(color)}">{num}</text>')
        out.append(f'<text x="{X+110}" y="{Y+70}" font-size="40" font-weight="700" fill="#FFFFFF">{escape(title)}</text>')

    def bullets(x, y, w, h, items, size=26, color=INK):
        lis = "".join(f'<li style="margin-bottom:12px">{escape(it)}</li>' for it in items)
        out.append(f'<foreignObject x="{px(x)}" y="{px(y)}" width="{px(w)}" height="{px(h)}">'
                   f'<div xmlns="http://www.w3.org/1999/xhtml" style="font-family:Arial,sans-serif;'
                   f'font-size:{size}px;color:{hx(color)};line-height:1.3">'
                   f'<ul style="margin:0;padding-left:{size+6}px">{lis}</ul></div></foreignObject>')

    def htmlbox(x, y, w, h, html):
        out.append(f'<foreignObject x="{px(x)}" y="{px(y)}" width="{px(w)}" height="{px(h)}">'
                   f'<div xmlns="http://www.w3.org/1999/xhtml" style="font-family:Arial,sans-serif">{html}</div></foreignObject>')

    cy0 = 5.5
    # ===== Column 1 =====
    card(C1, cy0, COLW, 8.0, 2, SEC_PROBLEM[0], SEC_PROBLEM[1])
    bullets(C1 + 0.45, cy0 + 1.45, COLW - 0.9, 6.3, SEC_PROBLEM[2])

    yB = 13.8
    card(C1, yB, COLW, 8.4, 3, SEC_RQ[0], SEC_RQ[1])
    htmlbox(C1 + 0.45, yB + 1.4, COLW - 0.9, 6.7,
        f'<div style="font-size:24px;color:{hx(BLUE)};font-weight:700;margin-bottom:4px">Research Question</div>'
        f'<div style="font-size:25px;color:{hx(INK)};line-height:1.3;margin-bottom:14px">{escape(RQ_TEXT)}</div>'
        f'<div style="background:{hx(LRED)};border:2px solid #F0A98C;border-radius:10px;padding:10px 12px;margin-bottom:10px">'
        f'<b style="color:{hx(AMBERD)};font-size:23px">H₀ (null)</b>'
        f'<div style="font-size:22px;color:{hx(INK)};line-height:1.25">{escape(H0_TEXT)}</div></div>'
        f'<div style="background:{hx(LTEAL)};border:2px solid #A7E0D6;border-radius:10px;padding:10px 12px">'
        f'<b style="color:{hx(TEAL)};font-size:23px">H₁ (alternative)</b>'
        f'<div style="font-size:22px;color:{hx(INK)};line-height:1.25">{escape(H1_TEXT)}</div></div>')

    yC = 22.5
    card(C1, yC, COLW, 10.4, 4, SEC_DATA[0], SEC_DATA[1])
    bullets(C1 + 0.45, yC + 1.45, COLW - 0.9, 8.6, SEC_DATA[2])

    # ===== Column 2 =====
    card(C2, cy0, COLW, 15.5, 5, SEC_APPROACH[0], SEC_APPROACH[1])
    bullets(C2 + 0.45, cy0 + 1.45, COLW - 0.9, 8.3, SEC_APPROACH[2])
    # high-level flow diagram
    dy = cy0 + 10.6
    def flowbox(x, w, t1, t2, fill, stroke):
        out.append(f'<rect x="{px(x)}" y="{px(dy)}" width="{px(w)}" height="{px(2.7)}" rx="10" fill="{hx(fill)}" stroke="{hx(stroke)}" stroke-width="2"/>')
        out.append(f'<text x="{px(x+w/2)}" y="{px(dy+1.15)}" text-anchor="middle" font-size="26" font-weight="700" fill="{hx(INK)}">{escape(t1)}</text>')
        if t2:
            out.append(f'<text x="{px(x+w/2)}" y="{px(dy+1.75)}" text-anchor="middle" font-size="21" fill="{hx(MUTE)}">{escape(t2)}</text>')
    def arrow(x):
        out.append(f'<path d="M{px(x)} {px(dy+1.35)} l{px(0.55)} 0 m-{px(0.18)} -{px(0.18)} l{px(0.18)} {px(0.18)} l-{px(0.18)} {px(0.18)}" stroke="{hx(MUTE)}" stroke-width="5" fill="none"/>')
    flowbox(C2 + 0.45, 3.9, "Predicted", "demand", LIGHT, BORDER)
    arrow(C2 + 4.5)
    flowbox(C2 + 5.2, 4.7, "Optimization Layer", "real-time + plan-ahead", LBLUE, BLUE)
    arrow(C2 + 10.05)
    flowbox(C2 + 10.7, 3.85, "Managed cluster", "overcommit, SLA met", LAMBER, AMBER)
    out.append(f'<text x="{px(C2+COLW/2)}" y="{px(dy+3.45)}" text-anchor="middle" font-size="21" fill="{hx(MUTE)}" font-style="italic">'
               f'Optimization guarantees feasibility and SLAs, unlike greedy rules.</text>')

    yE = 21.3
    card(C2, yE, COLW, 11.6, 6, SEC_KEY_TITLE := "Key Findings", BLUE)
    # KPI grid 2 cols x 3 (5 used)
    kx0, ky0 = C2 + 0.45, yE + 1.45
    kw, kh, kgx, kgy = 6.85, 2.7, 0.3, 0.3
    for i, (val, lab, col) in enumerate(KPIS):
        r_, c_ = divmod(i, 2)
        x = kx0 + c_ * (kw + kgx); y = ky0 + r_ * (kh + kgy)
        out.append(f'<rect x="{px(x)}" y="{px(y)}" width="{px(kw)}" height="{px(kh)}" rx="10" fill="{hx(LIGHT)}" stroke="{hx(BORDER)}" stroke-width="2"/>')
        out.append(f'<rect x="{px(x)}" y="{px(y)}" width="14" height="{px(kh)}" rx="6" fill="{hx(col)}"/>')
        out.append(f'<text x="{px(x+kw/2)}" y="{px(y+1.35)}" text-anchor="middle" font-size="58" font-weight="800" fill="{hx(col)}">{escape(val)}</text>')
        out.append(f'<text x="{px(x+kw/2)}" y="{px(y+2.15)}" text-anchor="middle" font-size="23" fill="{hx(INK)}">{escape(lab)}</text>')
    # one wide takeaway box in last slot
    x = kx0 + 1 * (kw + kgx); y = ky0 + 2 * (kh + kgy)
    out.append(f'<rect x="{px(x)}" y="{px(y)}" width="{px(kw)}" height="{px(kh)}" rx="10" fill="{hx(LBLUE)}" stroke="#A9C6EE" stroke-width="2"/>')
    htmlbox(x + 0.3, y + 0.25, kw - 0.6, kh - 0.4,
        f'<div style="font-size:22px;color:{hx(BLUE)};font-weight:700">Takeaway</div>'
        f'<div style="font-size:22px;color:{hx(INK)};line-height:1.25">Performance held as the cluster grew, keeping placements fast and SLAs intact.</div>')

    # ===== Column 3 =====
    card(C3, cy0, COLW, 15.5, 7, SEC_VIZ_TITLE := "Visualizations", BLUE)
    # --- bar chart ---
    out.append(f'<text x="{px(C3+0.5)}" y="{px(cy0+1.95)}" font-size="25" font-weight="700" fill="{hx(INK)}">Effective Memory Utilization (%)</text>')
    base_y = cy0 + 7.6; maxh = 4.7
    axL, axR = C3 + 1.0, C3 + COLW - 0.6
    out.append(f'<line x1="{px(axL)}" y1="{px(base_y)}" x2="{px(axR)}" y2="{px(base_y)}" stroke="{hx(MUTE)}" stroke-width="3"/>')
    bw_ = 3.2
    bxs = [C3 + 3.0, C3 + 8.4]
    bcols = [GREY, BLUE]
    for val, bxx, bcol, cat in zip(BAR_VALS, bxs, bcols, BAR_CATS):
        bhh = val / 100.0 * maxh
        out.append(f'<rect x="{px(bxx)}" y="{px(base_y-bhh)}" width="{px(bw_)}" height="{px(bhh)}" rx="6" fill="{hx(bcol)}"/>')
        out.append(f'<text x="{px(bxx+bw_/2)}" y="{px(base_y-bhh-0.2)}" text-anchor="middle" font-size="30" font-weight="800" fill="{hx(bcol)}">{val}%</text>')
        out.append(f'<text x="{px(bxx+bw_/2)}" y="{px(base_y+0.55)}" text-anchor="middle" font-size="22" fill="{hx(INK)}">{escape(cat)}</text>')
    out.append(f'<text x="{px(C3+0.5)}" y="{px(cy0+9.0)}" font-size="21" fill="{hx(MUTE)}" font-style="italic">Effective memory packed roughly doubles under the framework.</text>')
    # --- line chart ---
    ly0 = cy0 + 9.8
    out.append(f'<text x="{px(C3+0.5)}" y="{px(ly0+0.3)}" font-size="25" font-weight="700" fill="{hx(INK)}">Cluster Memory Over 24 Hours</text>')
    plL, plR = C3 + 1.0, C3 + COLW - 0.6
    plT, plB = ly0 + 0.6, ly0 + 4.0
    out.append(f'<rect x="{px(plL)}" y="{px(plT)}" width="{px(plR-plL)}" height="{px(plB-plT)}" fill="{hx(LIGHT)}" stroke="{hx(BORDER)}"/>')
    # SLA margin band (top 12%)
    out.append(f'<rect x="{px(plL)}" y="{px(plT)}" width="{px(plR-plL)}" height="{px((plB-plT)*0.12)}" fill="{hx(LRED)}" opacity="0.8"/>')
    n = len(LINE_HRS); span = (plR - plL) / (n - 1)
    def yval(v): return plB - (v / 100.0) * (plB - plT)
    # capacity dashed
    out.append(f'<line x1="{px(plL)}" y1="{px(yval(90))}" x2="{px(plR)}" y2="{px(yval(90))}" stroke="{hx(AMBERD)}" stroke-width="3" stroke-dasharray="10 7"/>')
    out.append(f'<text x="{px(plR-0.1)}" y="{px(yval(90)-0.12)}" text-anchor="end" font-size="19" fill="{hx(AMBERD)}">capacity</text>')
    pts = " ".join(f"{px(plL+i*span)},{px(yval(v))}" for i, v in enumerate(LINE_USE))
    out.append(f'<polyline points="{pts}" fill="none" stroke="{hx(BLUE)}" stroke-width="5"/>')
    for i, v in enumerate(LINE_USE):
        out.append(f'<circle cx="{px(plL+i*span)}" cy="{px(yval(v))}" r="7" fill="{hx(BLUE)}"/>')
    for i, hlab in enumerate(LINE_HRS):
        out.append(f'<text x="{px(plL+i*span)}" y="{px(plB+0.45)}" text-anchor="middle" font-size="19" fill="{hx(MUTE)}">{hlab}h</text>')
    out.append(f'<text x="{px(C3+0.5)}" y="{px(plB+1.0)}" font-size="21" fill="{hx(MUTE)}" font-style="italic">Usage stays under capacity with the SLA margin (red) kept free.</text>')

    yG = 21.3
    card(C3, yG, COLW, 5.8, 8, SEC_IMPL[0], SEC_IMPL[1])
    bullets(C3 + 0.45, yG + 1.45, COLW - 0.9, 4.0, SEC_IMPL[2], size=24)

    yH = 27.4
    card(C3, yH, COLW, 5.5, 9, SEC_LIMITS[0], SEC_LIMITS[1])
    bullets(C3 + 0.45, yH + 1.45, COLW - 0.9, 3.7, SEC_LIMITS[2], size=24)

    # ---- Footer ----
    fx, fy, fw, fh = 0.8, 33.1, 46.4, 2.1
    out.append(f'<rect x="{px(fx)}" y="{px(fy)}" width="{px(fw)}" height="{px(fh)}" rx="12" fill="{hx(NAVY)}"/>')
    out.append(f'<rect x="{px(fx)}" y="{px(fy)}" width="{px(fw/2)}" height="8" fill="{hx(AMBER)}"/>')
    out.append(f'<rect x="{px(fx+fw/2)}" y="{px(fy)}" width="{px(fw/2)}" height="8" fill="{hx(TEALB)}"/>')
    out.append(f'<text x="{px(fx+0.5)}" y="{px(fy+0.7)}" fill="#9FE7DA" font-size="24" font-weight="700">Selected References</text>')
    out.append(f'<text x="{px(fx+0.5)}" y="{px(fy+1.25)}" fill="#CBD8EE" font-size="20">Verma et al. (2015), Borg.  •  Google Cluster Usage Traces v3 (2019).</text>')
    out.append(f'<text x="{px(fx+0.5)}" y="{px(fy+1.7)}" fill="#CBD8EE" font-size="20">Maniezzo, Boschetti &amp; Stuetzle (2021), Matheuristics.</text>')
    out.append(f'<text x="{px(fx+19)}" y="{px(fy+0.7)}" fill="#9FE7DA" font-size="24" font-weight="700">Acknowledgement</text>')
    out.append(f'<text x="{px(fx+19)}" y="{px(fy+1.25)}" fill="#CBD8EE" font-size="20">Thanks to the prediction-layer team and the course instructor.</text>')
    out.append(f'<text x="{px(fx+34)}" y="{px(fy+0.7)}" fill="#9FE7DA" font-size="24" font-weight="700">Contact</text>')
    out.append(f'<text x="{px(fx+34)}" y="{px(fy+1.25)}" fill="#CBD8EE" font-size="20">Alrick Grandison — Optimization Layer</text>')

    out.append("</svg>")
    return "\n".join(out)


# ═════════════════════════════════════════════════════════════════════════════
#  PPTX RENDERER  (48 x 36 in; editable; export to PDF for printing)
# ═════════════════════════════════════════════════════════════════════════════
def build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    def C(h): return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width = Inches(PW)
    prs.slide_height = Inches(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes

    def rect(x, y, w, h, fill=None, line=None, lw=1.0, rounded=False, adj=0.04):
        shp = sh.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(x), Inches(y), Inches(w), Inches(h))
        if rounded:
            try: shp.adjustments[0] = adj
            except Exception: pass
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = C(line); shp.line.width = Pt(lw)
        shp.shadow.inherit = False
        return shp

    def text(shp, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
        """runs: list of paragraphs; each paragraph is list of (txt,size,bold,color,italic)."""
        tf = shp.text_frame; tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        for m in (tf.margin_left, ):
            pass
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        for i, para in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(6); p.space_before = Pt(0)
            for (txt, size, bold, color, italic) in para:
                r = p.add_run(); r.text = txt
                f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
                f.name = "Arial"; f.color.rgb = C(color)

    def textbox(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = sh.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        text(tb, runs, align=align, anchor=anchor)
        return tb

    def bullets_box(x, y, w, h, items, size=16, color=INK):
        tb = sh.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8); p.line_spacing = 1.05
            r = p.add_run(); r.text = "•  " + it
            f = r.font; f.size = Pt(size); f.name = "Arial"; f.color.rgb = C(color)
        return tb

    def card(x, y, w, h, num, title, color):
        rect(x, y, w, h, fill=WHITE, line=BORDER, lw=1.25, rounded=True, adj=0.03)
        rect(x, y, w, HH, fill=color, rounded=True, adj=0.06)
        rect(x, y + HH - 0.35, w, 0.35, fill=color)   # square off header bottom
        chip = sh.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.2),
                            Inches(0.7), Inches(0.7))
        chip.fill.solid(); chip.fill.fore_color.rgb = C(WHITE); chip.line.fill.background()
        chip.shadow.inherit = False
        text(chip, [[(str(num), 22, True, color, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb = sh.add_textbox(Inches(x + 1.1), Inches(y + 0.1), Inches(w - 1.3), Inches(HH - 0.2))
        text(tb, [[(title, 24, True, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE)

    # ---- Banner ----
    bx, by, bw, bh = 0.8, 0.6, 46.4, 4.6
    rect(bx, by, bw, bh, fill=NAVY, rounded=True, adj=0.03)
    rect(bx, by + bh - 0.16, bw / 2, 0.16, fill=AMBER)
    rect(bx + bw / 2, by + bh - 0.16, bw / 2, 0.16, fill=TEALB)
    rect(1.2, 1.2, 2.8, 2.8, fill=WHITE, rounded=True, adj=0.08)   # white chip behind logo
    if LOGO.exists():
        sh.add_picture(str(LOGO), Inches(1.45), Inches(1.45), Inches(2.3), Inches(2.3))
    if MASCOT.exists():
        sh.add_picture(str(MASCOT), Inches(43.6), Inches(0.85), Inches(3.1), Inches(3.1))
    textbox(bx, by + 0.55, bw, 1.1, [[(TITLE[0], 40, True, WHITE, False)]], align=PP_ALIGN.CENTER)
    textbox(bx, by + 1.6, bw, 0.9, [[(TITLE[1], 26, True, "9FE7DA", False)]], align=PP_ALIGN.CENTER)
    textbox(bx, by + 2.55, bw, 0.6, [[(AUTHOR, 20, False, "CBD8EE", False)]], align=PP_ALIGN.CENTER)
    textbox(bx, by + 3.2, bw, 0.6,
            [[(f"{COURSE}    |    {INSTR}    |    {INSTN}    |    {DATE}", 17, False, "9FB4D6", False)]],
            align=PP_ALIGN.CENTER)

    cy0 = 5.5
    # ===== Column 1 =====
    card(C1, cy0, COLW, 8.0, 2, SEC_PROBLEM[0], SEC_PROBLEM[1])
    bullets_box(C1 + 0.4, cy0 + 1.35, COLW - 0.8, 6.3, SEC_PROBLEM[2], size=17)

    yB = 13.8
    card(C1, yB, COLW, 8.4, 3, SEC_RQ[0], SEC_RQ[1])
    textbox(C1 + 0.4, yB + 1.3, COLW - 0.8, 1.8, [
        [("Research Question", 16, True, BLUE, False)],
        [(RQ_TEXT, 17, False, INK, False)],
    ])
    rb = rect(C1 + 0.4, yB + 3.5, COLW - 0.8, 2.0, fill=LRED, line="F0A98C", lw=1.0, rounded=True, adj=0.04)
    text(rb, [[("H₀ (null)  ", 15, True, AMBERD, False), (H0_TEXT, 15, False, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    gb = rect(C1 + 0.4, yB + 5.7, COLW - 0.8, 2.2, fill=LTEAL, line="A7E0D6", lw=1.0, rounded=True, adj=0.04)
    text(gb, [[("H₁ (alt.)  ", 15, True, TEAL, False), (H1_TEXT, 15, False, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)

    yC = 22.5
    card(C1, yC, COLW, 10.4, 4, SEC_DATA[0], SEC_DATA[1])
    bullets_box(C1 + 0.4, yC + 1.35, COLW - 0.8, 8.6, SEC_DATA[2], size=17)

    # ===== Column 2 =====
    card(C2, cy0, COLW, 15.5, 5, SEC_APPROACH[0], SEC_APPROACH[1])
    bullets_box(C2 + 0.4, cy0 + 1.35, COLW - 0.8, 8.6, SEC_APPROACH[2], size=17)
    dy = cy0 + 10.6
    def flow(x, w, t1, t2, fill, line):
        b = rect(x, dy, w, 2.5, fill=fill, line=line, lw=1.25, rounded=True, adj=0.05)
        text(b, [[(t1, 17, True, INK, False)]] + ([[(t2, 13, False, MUTE, False)]] if t2 else []),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    def arr(x, w):
        a = sh.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(dy + 1.0), Inches(w), Inches(0.5))
        a.fill.solid(); a.fill.fore_color.rgb = C(MUTE); a.line.fill.background(); a.shadow.inherit = False
    flow(C2 + 0.4, 3.9, "Predicted demand", "", LIGHT, BORDER)
    arr(C2 + 4.4, 0.7)
    flow(C2 + 5.2, 4.7, "Optimization Layer", "real-time + plan-ahead", LBLUE, BLUE)
    arr(C2 + 10.0, 0.7)
    flow(C2 + 10.75, 3.85, "Managed cluster", "overcommit, SLA met", LAMBER, AMBER)
    textbox(C2 + 0.4, dy + 2.6, COLW - 0.8, 0.6,
            [[("Optimization guarantees feasibility and SLAs, unlike greedy rules.", 13, False, MUTE, True)]],
            align=PP_ALIGN.CENTER)

    yE = 21.3
    card(C2, yE, COLW, 11.6, 6, "Key Findings", BLUE)
    kx0, ky0 = C2 + 0.4, yE + 1.35
    kw, kh, kgx, kgy = 6.9, 2.7, 0.3, 0.3
    for i, (val, lab, col) in enumerate(KPIS):
        r_, c_ = divmod(i, 2)
        x = kx0 + c_ * (kw + kgx); y = ky0 + r_ * (kh + kgy)
        rect(x, y, kw, kh, fill=LIGHT, line=BORDER, lw=1.0, rounded=True, adj=0.05)
        rect(x, y, 0.12, kh, fill=col)
        textbox(x, y + 0.25, kw, 1.6, [[(val, 40, True, col, False)]], align=PP_ALIGN.CENTER)
        textbox(x, y + 1.85, kw, 0.7, [[(lab, 15, False, INK, False)]], align=PP_ALIGN.CENTER)
    x = kx0 + (kw + kgx); y = ky0 + 2 * (kh + kgy)
    tk = rect(x, y, kw, kh, fill=LBLUE, line="A9C6EE", lw=1.0, rounded=True, adj=0.05)
    text(tk, [[("Takeaway  ", 15, True, BLUE, False)],
              [("Performance held as the cluster grew — fast placements, SLAs intact.", 14, False, INK, False)]],
         anchor=MSO_ANCHOR.MIDDLE)

    # ===== Column 3 =====
    card(C3, cy0, COLW, 15.5, 7, "Visualizations", BLUE)
    # bar chart
    cd1 = CategoryChartData(); cd1.categories = BAR_CATS
    cd1.add_series("Memory used (%)", BAR_VALS)
    gf1 = sh.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                       Inches(C3 + 0.5), Inches(cy0 + 1.4), Inches(COLW - 1.0), Inches(5.6), cd1)
    ch1 = gf1.chart; ch1.has_legend = False
    ch1.has_title = True; ch1.chart_title.text_frame.text = "Effective Memory Utilization (%)"
    for p in ch1.chart_title.text_frame.paragraphs:
        for r in p.runs: r.font.size = Pt(16); r.font.name = "Arial"; r.font.bold = True
    plot1 = ch1.plots[0]; plot1.has_data_labels = True
    plot1.data_labels.number_format = '0"%"'; plot1.data_labels.number_format_is_linked = False
    plot1.data_labels.font.size = Pt(16); plot1.data_labels.font.bold = True
    ser = plot1.series[0]
    ser.points[0].format.fill.solid(); ser.points[0].format.fill.fore_color.rgb = C(GREY)
    ser.points[1].format.fill.solid(); ser.points[1].format.fill.fore_color.rgb = C(BLUE)
    ch1.category_axis.tick_labels.font.size = Pt(13)
    ch1.value_axis.tick_labels.font.size = Pt(11); ch1.value_axis.maximum_scale = 100
    textbox(C3 + 0.5, cy0 + 7.0, COLW - 1.0, 0.5,
            [[("Effective memory packed roughly doubles under the framework.", 13, False, MUTE, True)]])
    # line chart
    cd2 = CategoryChartData(); cd2.categories = [f"{h}h" for h in LINE_HRS]
    cd2.add_series("Usage (%)", LINE_USE)
    cd2.add_series("Capacity", LINE_CAP)
    gf2 = sh.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                       Inches(C3 + 0.5), Inches(cy0 + 7.6), Inches(COLW - 1.0), Inches(6.2), cd2)
    ch2 = gf2.chart; ch2.has_legend = True
    from pptx.enum.chart import XL_LEGEND_POSITION
    ch2.legend.position = XL_LEGEND_POSITION.BOTTOM; ch2.legend.include_in_layout = False
    ch2.legend.font.size = Pt(12)
    ch2.has_title = True; ch2.chart_title.text_frame.text = "Cluster Memory Over 24 Hours"
    for p in ch2.chart_title.text_frame.paragraphs:
        for r in p.runs: r.font.size = Pt(16); r.font.name = "Arial"; r.font.bold = True
    s0 = ch2.plots[0].series[0]; s0.format.line.color.rgb = C(BLUE); s0.format.line.width = Pt(2.5)
    s1 = ch2.plots[0].series[1]; s1.format.line.color.rgb = C(AMBERD); s1.format.line.width = Pt(2.0)
    ch2.category_axis.tick_labels.font.size = Pt(12)
    ch2.value_axis.tick_labels.font.size = Pt(11); ch2.value_axis.maximum_scale = 100

    yG = 21.3
    card(C3, yG, COLW, 5.8, 8, SEC_IMPL[0], SEC_IMPL[1])
    bullets_box(C3 + 0.4, yG + 1.35, COLW - 0.8, 4.0, SEC_IMPL[2], size=16)

    yH = 27.4
    card(C3, yH, COLW, 5.5, 9, SEC_LIMITS[0], SEC_LIMITS[1])
    bullets_box(C3 + 0.4, yH + 1.35, COLW - 0.8, 3.7, SEC_LIMITS[2], size=16)

    # ---- Footer ----
    fx, fy, fw, fh = 0.8, 33.1, 46.4, 2.1
    rect(fx, fy, fw, fh, fill=NAVY, rounded=True, adj=0.05)
    rect(fx, fy, fw / 2, 0.08, fill=AMBER); rect(fx + fw / 2, fy, fw / 2, 0.08, fill=TEALB)
    textbox(fx + 0.5, fy + 0.25, 18, fh - 0.4, [
        [("Selected References", 15, True, "9FE7DA", False)],
        [("Verma et al. (2015), Borg.  •  Google Cluster Traces v3 (2019).", 12, False, "CBD8EE", False)],
        [("Maniezzo, Boschetti & Stuetzle (2021), Matheuristics.", 12, False, "CBD8EE", False)],
    ])
    textbox(fx + 19, fy + 0.25, 14, fh - 0.4, [
        [("Acknowledgement", 15, True, "9FE7DA", False)],
        [("Thanks to the prediction-layer team and the", 12, False, "CBD8EE", False)],
        [("course instructor for their guidance.", 12, False, "CBD8EE", False)],
    ])
    textbox(fx + 34, fy + 0.25, 12, fh - 0.4, [
        [("Contact", 15, True, "9FE7DA", False)],
        [("Alrick Grandison", 12, False, "CBD8EE", False)],
        [("Optimization Layer", 12, False, "CBD8EE", False)],
    ])

    prs.save(str(path))


def main():
    svg_path = HERE / "capstone_poster.svg"
    svg_path.write_text(build_svg(), encoding="utf-8")
    print(f"SVG  saved: {svg_path}")
    pptx_path = HERE / "capstone_poster.pptx"
    build_pptx(pptx_path)
    print(f"PPTX saved: {pptx_path}")


if __name__ == "__main__":
    main()
