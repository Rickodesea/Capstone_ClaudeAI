from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.oxml.ns as ns
from lxml import etree

# ── Color Palette ──────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)
C_GOLD      = RGBColor(0xFF, 0xD1, 0x66)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCA, 0xE9, 0xFF)
C_DARK_BOX  = RGBColor(0x1A, 0x2E, 0x44)
C_GREEN     = RGBColor(0x4C, 0xAF, 0x50)
C_RED       = RGBColor(0xEF, 0x53, 0x50)
C_YELLOW    = RGBColor(0xFF, 0xEE, 0x58)
C_GRAY      = RGBColor(0x78, 0x90, 0x9C)
C_BRONZE    = RGBColor(0xE0, 0x87, 0x6A)
C_SILVER    = RGBColor(0xA8, 0xDA, 0xDC)
C_TEAL      = RGBColor(0x00, 0x96, 0x88)
C_PURPLE    = RGBColor(0x7C, 0x4D, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    return s

def rect(s, l, t, w, h, fill=None, border=None, bw=Pt(1.5)):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = bw
    else:
        sh.line.fill.background()
    return sh

def tb(s, text, l, t, w, h, sz=Pt(13), bold=False, col=C_WHITE, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold; r.font.color.rgb = col
    return box

def bullets(s, items, l, t, w, h, sz=Pt(11.5), col=C_WHITE, prefix="▸  ", spacing=Pt(6)):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = spacing
        r = p.add_run(); r.text = prefix + item
        r.font.size = sz; r.font.color.rgb = col
    return box

def header(s, title, sub=None):
    rect(s, 0, 0, 0.08, 7.5, fill=C_ACCENT)
    tb(s, title, 0.2, 0.18, 12.8, 0.65, sz=Pt(26), bold=True)
    if sub:
        tb(s, sub, 0.2, 0.82, 12.8, 0.42, sz=Pt(14), col=C_ACCENT)

def arrow(s, x1, y1, x2, y2, col=C_ACCENT, w=Pt(2)):
    c = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = col; c.line.width = w
    ln = c.line._ln
    etree.SubElement(ln, ns.qn('a:tailEnd')).set('type', 'none')
    h = etree.SubElement(ln, ns.qn('a:headEnd'))
    h.set('type', 'arrow'); h.set('w', 'med'); h.set('len', 'med')

def label_box(s, text, l, t, w, h, fill=C_DARK_BOX, border=C_ACCENT, sz=Pt(11.5), bold=False, col=C_WHITE):
    rect(s, l, t, w, h, fill=fill, border=border)
    box = s.shapes.add_textbox(Inches(l+0.08), Inches(t+0.06), Inches(w-0.16), Inches(h-0.1))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold; r.font.color.rgb = col

def divider(s, y, col=C_ACCENT, thickness=Pt(1)):
    rect(s, 0.2, y, 12.9, 0.03, fill=col)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
# Top accent bar
rect(s, 0, 0, 13.33, 0.12, fill=C_ACCENT)
rect(s, 0, 7.38, 13.33, 0.12, fill=C_ACCENT)

tb(s, "Smart Scheduling for Multi-Tenant Kubernetes Clusters",
   0.7, 0.8, 11.9, 1.6, sz=Pt(36), bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Predictive Admission  ·  Fair Resource Allocation  ·  SLA Compliance",
   0.7, 2.35, 11.9, 0.55, sz=Pt(16), col=C_ACCENT, align=PP_ALIGN.CENTER)

divider(s, 3.05)

# Two-column summary boxes
for i, (hdr, body, col) in enumerate([
    ("THE PROBLEM", "Cloud clusters waste 40–60% of memory because schedulers trust declared resource requests, not actual usage. SLA violations happen after the fact. Fairness across tenants is accidental.", C_RED),
    ("OUR APPROACH", "A Kubernetes scheduler plugin that predicts real memory demand, places jobs fairly using Dominant Resource Fairness, and enforces SLAs at runtime via dynamic cgroup controls.", C_GREEN),
]):
    x = 0.4 + i * 6.5
    rect(s, x, 3.25, 6.1, 0.45, fill=col)
    tb(s, hdr, x, 3.25, 6.1, 0.45, sz=Pt(13), bold=True,
       col=C_BG if col == C_GREEN else C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, 3.7, 6.1, 1.9, fill=C_DARK_BOX, border=col)
    tb(s, body, x+0.12, 3.76, 5.88, 1.78, sz=Pt(12), col=C_LIGHT)

rect(s, 0.4, 5.8, 12.5, 0.55, fill=RGBColor(0x00, 0x2A, 0x3A), border=C_GOLD)
tb(s, "DAMO 699 Capstone  |  Group 3  |  Supervisor: Hany Osman  |  University of Niagara Falls  |  Spring 2026",
   0.5, 5.82, 12.3, 0.5, sz=Pt(12), col=C_GOLD, align=PP_ALIGN.CENTER)
tb(s, "Team: Tha Pyay Hmu  ·  Lhagii Tsogtbayar  ·  Nadia Ríos  ·  Jorge Mendoza  ·  Alrick Grandison",
   0.5, 6.45, 12.3, 0.4, sz=Pt(11), col=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "The Problem", "Three gaps the default Kubernetes scheduler cannot close")

# Stat banner
rect(s, 0.2, 1.35, 12.9, 0.72, fill=RGBColor(0x00, 0x2A, 0x3A), border=C_GOLD)
for i, (num, lbl) in enumerate([
    ("40–60%", "avg cluster memory\nutilization (wasted)"),
    ("< 80%", "SLA compliance under\nhigh tenant load"),
    ("0%", "native fairness in\ndefault K8s scheduler"),
    ("70%", "AI training perf lost to\npoor node placement"),
]):
    x = 0.5 + i * 3.2
    tb(s, num, x, 1.38, 2.8, 0.38, sz=Pt(20), bold=True, col=C_GOLD, align=PP_ALIGN.CENTER)
    tb(s, lbl, x, 1.73, 2.8, 0.3, sz=Pt(9), col=C_LIGHT, align=PP_ALIGN.CENTER)

# Three problem boxes
for i, (title, body, col) in enumerate([
    ("No Prediction",
     "The scheduler looks at declared resource requests — not actual usage.\n\n"
     "Tenants over-declare by 30–50% as a safety buffer. The scheduler treats those inflated numbers as real, keeping nodes 'full' when they are actually half-empty.\n\n"
     "Result: hardware sits idle while jobs queue. Memory OOM kills occur when peaks aren't anticipated.",
     C_RED),
    ("No Fairness",
     "Kubernetes processes job submissions in arrival order — first come, first served.\n\n"
     "A tenant who submits 100 batch jobs before another tenant submits one critical job forces that critical job to wait.\n\n"
     "Result: some tenants monopolize resources; others starve. No mechanism prevents this.",
     C_GOLD),
    ("No SLA Enforcement",
     "SLA violations are discovered after they occur — an OOM kill, a latency spike, a job that missed its deadline.\n\n"
     "There is no admission gate that asks: 'Will admitting this job cause a co-tenant's SLA to be violated in 30 minutes?'\n\n"
     "Result: reactive damage control instead of proactive prevention.",
     C_ACCENT),
]):
    x = 0.2 + i * 4.35
    rect(s, x, 2.25, 4.1, 0.5, fill=col)
    tb(s, title, x, 2.25, 4.1, 0.5, sz=Pt(14), bold=True,
       col=C_BG if col != C_ACCENT else C_BG, align=PP_ALIGN.CENTER)
    rect(s, x, 2.75, 4.1, 3.9, fill=C_DARK_BOX, border=col)
    tb(s, body, x+0.12, 2.82, 3.88, 3.78, sz=Pt(11.5))

rect(s, 0.2, 6.82, 12.9, 0.45, fill=RGBColor(0x00, 0x3A, 0x52), border=C_GREEN)
tb(s, "Our project directly addresses all three gaps with a single integrated Kubernetes scheduler plugin.",
   0.35, 6.84, 12.6, 0.4, sz=Pt(13), bold=True, col=C_GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR APPROACH
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Our Approach", "Three integrated components — one Kubernetes scheduler plugin")

# What we are NOT building (quick note)
rect(s, 0.2, 1.3, 12.9, 0.38, fill=RGBColor(0x1A, 0x1A, 0x1A), border=C_GRAY, bw=Pt(1))
tb(s, "We are not building a cloud platform like AWS. We are building the scheduling decision model that runs inside a Kubernetes cluster, and evaluating it via simulation on the Google Cluster Trace v3.",
   0.35, 1.32, 12.6, 0.34, sz=Pt(11), col=C_GRAY)

# Three component boxes
comp_data = [
    ("1. Predictive Admission\nControl", C_GREEN,
     "What it does:\nBefore scheduling a job, predict what the node's total memory will look like if we admit it — not based on declared requests, but on actual historical usage per tenant.\n\nHow:\nRandom Forest model trained on Google Cluster Trace. Features include tenant's historical request-to-actual ratio, job type, time of day. Predicts P95 memory peak.\n\nResult:\nOnly admit if predicted combined memory stays below threshold. Enables safe overcommitment."),
    ("2. Fairness-Aware\nPlacement (DRF)", C_GOLD,
     "What it does:\nAmong all nodes that pass the admission check, rank them by which placement is fairest to all tenants currently running in the cluster.\n\nHow:\nDominant Resource Fairness (DRF): find which resource (CPU or memory) each tenant consumes the most of. Give priority to the tenant whose dominant resource share is smallest — the most underserved tenant goes next.\n\nResult:\nNo tenant monopolizes. Fairness is mathematical and provable, not accidental."),
    ("3. Runtime SLA\nEnforcement", C_ACCENT,
     "What it does:\nOnce a job is running, continuously monitor its memory and CPU usage via Prometheus. Enforce that critical (SLA-bound) jobs always get their resources, even if batch jobs need to be throttled.\n\nHow:\nKubernetes QoS classes + dynamic cgroup weight adjustments. Critical jobs = Guaranteed QoS (never throttled). Batch jobs = Burstable (yield under pressure).\n\nResult:\nSLA violations prevented at runtime, not just at admission time."),
]
for i, (title, col, body) in enumerate(comp_data):
    x = 0.2 + i * 4.35
    rect(s, x, 1.85, 4.1, 0.72, fill=col)
    tb(s, title, x, 1.85, 4.1, 0.72, sz=Pt(13), bold=True,
       col=C_BG, align=PP_ALIGN.CENTER)
    rect(s, x, 2.57, 4.1, 4.3, fill=C_DARK_BOX, border=col)
    tb(s, body, x+0.12, 2.63, 3.88, 4.18, sz=Pt(10.5))

rect(s, 0.2, 7.05, 12.9, 0.3, fill=RGBColor(0x00, 0x3A, 0x52), border=C_GOLD)
tb(s, "Novel combination: no existing paper combines all three. Chaudhari (2025) calls for this exact system — we build it.",
   0.35, 7.07, 12.6, 0.26, sz=Pt(11), col=C_GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "System Architecture", "Job submission → Admission → Scheduling → Runtime enforcement")

# Pipeline: horizontal flow
cx_start = 0.3
box_w = 2.55
box_h = 1.3
gap = 0.22
y_top = 1.5

stages = [
    ("Job\nSubmission", C_GRAY,
     "Tenant submits pod via kubectl. Declares CPU + memory requests."),
    ("Admission\nWebhook", C_GREEN,
     "ML model predicts actual peak memory. Accept or queue."),
    ("Scheduler\nPlugin", C_GOLD,
     "Filter → DRF Score → Temporal co-location score → Bind to node."),
    ("Runtime\nController", C_ACCENT,
     "Prometheus monitors live usage. Adjusts cgroup weights every 5s."),
    ("Outcome", C_GREEN,
     "Job runs safely. SLA maintained. Tenant share is fair."),
]
for i, (title, col, desc) in enumerate(stages):
    x = cx_start + i * (box_w + gap)
    rect(s, x, y_top, box_w, 0.5, fill=col)
    tb(s, title, x, y_top, box_w, 0.5, sz=Pt(12), bold=True,
       col=C_BG, align=PP_ALIGN.CENTER)
    rect(s, x, y_top + 0.5, box_w, box_h - 0.5, fill=C_DARK_BOX, border=col)
    tb(s, desc, x+0.1, y_top + 0.55, box_w-0.2, box_h-0.6, sz=Pt(10.5))
    if i < len(stages) - 1:
        ax = x + box_w + 0.01
        arrow(s, ax, y_top + 0.65, ax + gap + 0.01, y_top + 0.65, col=col, w=Pt(2))

# Kubernetes layers diagram below
rect(s, 0.2, 3.1, 12.9, 0.4, fill=RGBColor(0x00, 0x3A, 0x52), border=C_ACCENT)
tb(s, "Kubernetes Scheduling Framework Hook Points",
   0.35, 3.12, 12.6, 0.36, sz=Pt(12), bold=True, col=C_ACCENT)

hook_data = [
    ("Admission Webhook\n(before scheduler)", C_GREEN,
     "▸ Fetch tenant history\n▸ Run Random Forest → predict P95 memory\n▸ Compare vs. threshold\n▸ ADMIT or QUEUE"),
    ("PreFilter + Filter\n(scheduler phase 1)", C_GRAY,
     "▸ Enforce namespace ResourceQuotas\n▸ Eliminate nodes with insufficient physical memory\n▸ Check SloPolicy CRD constraints"),
    ("Score\n(scheduler phase 2)", C_GOLD,
     "▸ DRF fairness score per node\n▸ Temporal co-location score\n▸ (lower dominant share tenant = higher priority)"),
    ("Bind + Runtime\n(post-placement)", C_ACCENT,
     "▸ Assign pod to highest-scoring node\n▸ Set Guaranteed/Burstable QoS class\n▸ Start Prometheus monitoring + cgroup controller"),
]
for i, (title, col, body) in enumerate(hook_data):
    x = 0.2 + i * 3.28
    rect(s, x, 3.6, 3.1, 0.45, fill=col)
    tb(s, title, x, 3.6, 3.1, 0.45, sz=Pt(10), bold=True,
       col=C_BG, align=PP_ALIGN.CENTER)
    rect(s, x, 4.05, 3.1, 2.2, fill=C_DARK_BOX, border=col)
    tb(s, body, x+0.1, 4.1, 2.9, 2.1, sz=Pt(10))

# Memory bar at bottom
tb(s, "Memory safety model:", 0.2, 6.4, 3.5, 0.3, sz=Pt(11), bold=True, col=C_LIGHT)
rect(s, 0.2, 6.72, 12.9, 0.52, fill=RGBColor(0x12, 0x26, 0x38), border=C_ACCENT)
rect(s, 0.2, 6.72, 5.5, 0.52, fill=C_TEAL)
tb(s, "Current committed ~45%", 0.3, 6.72, 5.3, 0.52, sz=Pt(9.5), col=C_BG, align=PP_ALIGN.CENTER)
rect(s, 5.7, 6.72, 5.2, 0.52, fill=RGBColor(0xCC, 0xA0, 0x00))
tb(s, "Predictive headroom (safe overcommit ~40%)", 5.8, 6.72, 5.0, 0.52, sz=Pt(9.5), col=C_BG, align=PP_ALIGN.CENTER)
rect(s, 10.9, 6.72, 2.2, 0.52, fill=C_DARK_BOX, border=C_GRAY)
tb(s, "Safety buffer", 10.95, 6.72, 2.1, 0.52, sz=Pt(9.5), col=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — OPTIMIZATION MODEL DESIGN
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Suggested Optimization Model Design",
       "Kovalenko (2024) base structure → extended with prediction, DRF fairness, and SLA constraints")

# ── Left column: what we keep vs. change ──────────────────────────────────────
rect(s, 0.2, 1.35, 4.0, 0.4, fill=C_TEAL)
tb(s, "Kovalenko Base  →  Our Extensions", 0.2, 1.35, 4.0, 0.4,
   sz=Pt(11), bold=True, col=C_BG, align=PP_ALIGN.CENTER)
rect(s, 0.2, 1.75, 4.0, 5.45, fill=C_DARK_BOX, border=C_TEAL)

keep_change = [
    ("KEEP", C_GREEN, [
        "Binary pod→node assignment  x_{n,j} ∈ {0,1}",
        "Per-node CPU + memory capacity constraints",
        "Multi-objective structure",
        "Discrete time horizon T",
    ]),
    ("REMOVE", C_RED, [
        "Server power on/off objective",
        "  (K8s nodes are pre-provisioned VMs —",
        "   power mgmt is cloud provider's job)",
    ]),
    ("ADD / REPLACE", C_GOLD, [
        "Declared requests  →  RF-predicted m̂_j(t), ĉ_j(t)",
        "DRF dominant share constraint  d_k ≤ 1/|K| + ε",
        "SLA deadline  T_start + dur_j ≤ D_j",
        "Peak-time safety check at predicted t*_j",
    ]),
]
y_inner = 1.82
for tag, col, items in keep_change:
    tb(s, tag, 0.32, y_inner, 1.0, 0.28, sz=Pt(9.5), bold=True, col=col)
    y_inner += 0.28
    for item in items:
        tb(s, "  " + item, 0.32, y_inner, 3.78, 0.26, sz=Pt(9.5), col=C_LIGHT)
        y_inner += 0.26
    y_inner += 0.08

# ── Middle column: objective + variable definitions ────────────────────────────
rect(s, 4.4, 1.35, 4.55, 0.4, fill=C_GOLD)
tb(s, "Objective Function + Key Variables", 4.4, 1.35, 4.55, 0.4,
   sz=Pt(11), bold=True, col=C_BG, align=PP_ALIGN.CENTER)
rect(s, 4.4, 1.75, 4.55, 5.45, fill=C_DARK_BOX, border=C_GOLD)

obj_lines = [
    ("maximize:", C_GOLD),
    ("  w₁ · avg memory utilization", C_YELLOW),
    ("  − w₂ · SLA violation penalty", C_YELLOW),
    ("  − w₃ · dominant share spread", C_YELLOW),
    ("", C_WHITE),
    ("Suggested weights:", C_LIGHT),
    ("  w₁ = 0.50  (utilization)", C_WHITE),
    ("  w₂ = 0.35  (SLA protection)", C_WHITE),
    ("  w₃ = 0.15  (fairness)", C_WHITE),
    ("", C_WHITE),
    ("Key Variables:", C_GOLD),
    ("  x_{n,j} ∈ {0,1}   pod j → node n", C_LIGHT),
    ("  q_j ∈ {0,1}       job queued", C_LIGHT),
    ("  d_k ∈ [0,1]       tenant dominant share", C_LIGHT),
    ("  m̂_j(t)            RF-predicted memory", C_GREEN),
    ("  ĉ_j(t)            RF-predicted CPU", C_GREEN),
    ("  dur_j             RF-predicted duration", C_GREEN),
    ("", C_WHITE),
    ("Implementation note:", C_GOLD),
    ("  Full model is NP-hard. Use greedy", C_LIGHT),
    ("  sequential admission in simulation:", C_LIGHT),
    ("  C1→C2→C4/C5→C6 → admit or queue", C_LIGHT),
]
box = s.shapes.add_textbox(Inches(4.52), Inches(1.82), Inches(4.3), Inches(5.28))
box.word_wrap = True; tf = box.text_frame; tf.word_wrap = True
for i, (line, col) in enumerate(obj_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(1.5)
    r = p.add_run(); r.text = line
    r.font.size = Pt(10); r.font.color.rgb = col

# ── Right column: constraints ─────────────────────────────────────────────────
rect(s, 9.15, 1.35, 4.0, 0.4, fill=C_ACCENT)
tb(s, "Constraints  (C1 – C8)", 9.15, 1.35, 4.0, 0.4,
   sz=Pt(11), bold=True, col=C_BG, align=PP_ALIGN.CENTER)
rect(s, 9.15, 1.75, 4.0, 5.45, fill=C_DARK_BOX, border=C_ACCENT)

constraints = [
    ("C1", "Memory at predicted peak:", C_RED),
    ("",   "Σ_j x_{n,j}·m̂_j(t*) ≤ α·RAM_n  ∀n", C_YELLOW),
    ("C2", "CPU per time period:", C_BRONZE),
    ("",   "Σ_j x_{n,j}·ĉ_j(t) ≤ α·CPU_n  ∀n,t", C_YELLOW),
    ("C3", "Single assignment:", C_ACCENT),
    ("",   "Σ_n x_{n,j} + q_j = 1  ∀j", C_YELLOW),
    ("C4", "DRF dominant share:", C_GOLD),
    ("",   "d_k = max(mem_share_k, cpu_share_k)", C_YELLOW),
    ("C5", "Fairness bound:", C_GOLD),
    ("",   "d_k ≤ (1/|K|) + ε  ∀k", C_YELLOW),
    ("C6", "SLA deadline:", C_GREEN),
    ("",   "T_start + dur_j ≤ D_j  (if placed)", C_YELLOW),
    ("C7", "Temporal co-location:", C_SILVER),
    ("",   "Σ_j x_{n,j}·m̂_j(t) ≤ α·RAM_n  ∀n,t", C_YELLOW),
    ("C8", "Binary + non-negativity", C_GRAY),
]
box2 = s.shapes.add_textbox(Inches(9.25), Inches(1.82), Inches(3.8), Inches(5.28))
box2.word_wrap = True; tf2 = box2.text_frame; tf2.word_wrap = True
for i, (tag, line, col) in enumerate(constraints):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(2)
    if tag:
        r = p.add_run(); r.text = tag + "  "
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = col
        r2 = p.add_run(); r2.text = line
        r2.font.size = Pt(10); r2.font.color.rgb = C_LIGHT
    else:
        r = p.add_run(); r.text = "      " + line
        r.font.size = Pt(9.5); r.font.color.rgb = C_YELLOW

# ── Bottom bar: novel contributions ───────────────────────────────────────────
rect(s, 0.2, 7.25, 12.9, 0.08, fill=C_ACCENT)
rect(s, 0.2, 7.05, 12.9, 0.2, fill=RGBColor(0x00, 0x2A, 0x3A))
tb(s,
   "Novel: no existing K8s paper combines predicted utilization in constraints (C1,C7) + DRF fairness (C4,C5) + SLA deadlines (C6) in one model.  α = safety ratio (e.g., 0.90).",
   0.35, 7.06, 12.6, 0.18, sz=Pt(10), col=C_GOLD, align=PP_ALIGN.CENTER)

# ── Prediction model sub-box ────────────────────────────────────────────────
rect(s, 0.2, 6.38, 12.9, 0.48, fill=RGBColor(0x0D, 0x2A, 0x1A), border=C_GREEN)
tb(s, "Prediction inputs (feeds constraints above):", 0.32, 6.4, 3.5, 0.22,
   sz=Pt(9.5), bold=True, col=C_GREEN)
tb(s,
   "Random Forest → m̂_j(t), ĉ_j(t), dur_j  |  "
   "Features: tenant hist. ratio · declared request · job class · hour-of-day  |  "
   "Pulse shape: 10% ramp-up → 80% plateau → 10% ramp-down  |  "
   "Preprocessing: Savitzky-Golay + min-max (Kofi 2025)",
   0.32, 6.6, 12.6, 0.22, sz=Pt(9.5), col=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RESEARCH FOUNDATION (was SLIDE 5)
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Research Foundation", "18 papers reviewed — each maps to a specific system component")

# Component map
comp_map = [
    ("Prediction Layer", C_GREEN, [
        "Resource Central (Cortez 2017) — RF model, tenant history features",
        "Doukha (2025) — RF beats LSTM; MAPE 2.65% justifies our model choice",
        "Kofi (2025) — Google Cluster Trace preprocessing; R²=0.99 benchmark",
        "Wang & Yang (2025) — K8s LSTM+DQN baseline to compare against",
    ]),
    ("Fairness & Placement", C_GOLD, [
        "DRF (Ghodsi 2011) — formal fairness algorithm, proven properties",
        "Coach / Reidys (2025) — temporal co-location: non-overlapping peaks",
        "Alatawi (2025) — Gini coefficient as second fairness metric",
        "Kovalenko & Zhdanova (2024) — math optimization formalization",
    ]),
    ("SLA & Enforcement", C_ACCENT, [
        "Priya (2025) — SloPolicy CRD + Kubernetes plugin architecture",
        "Blagodurov — cgroups dynamic weights; critical vs batch co-location",
        "Liu & Guitart (2025) — in-node DRC for group-aware cgroup tuning",
        "Zhao et al. (2021) — formal admission control with SLA constraints",
    ]),
    ("Evaluation & Context", C_BRONZE, [
        "Chaudhari (2025) — gap anchor; the unbuilt framework we construct",
        "Jiang Zhi (2025) — Clovers simulator + real trace datasets",
        "Perera (2025) — RL risks; model drift; justifies our RF approach",
        "Atropos / Hu (2025) — overload fallback: target culprit, not victim",
    ]),
]
for i, (title, col, items) in enumerate(comp_map):
    x = 0.2 + i * 3.28
    rect(s, x, 1.38, 3.1, 0.45, fill=col)
    tb(s, title, x, 1.38, 3.1, 0.45, sz=Pt(11.5), bold=True, col=C_BG, align=PP_ALIGN.CENTER)
    rect(s, x, 1.83, 3.1, 4.35, fill=C_DARK_BOX, border=col)
    bullets(s, items, x+0.08, 1.9, 2.95, 4.2, sz=Pt(10.5), col=C_LIGHT)

# Gap statement
rect(s, 0.2, 6.32, 12.9, 0.52, fill=RGBColor(0x00, 0x2A, 0x3A), border=C_GOLD, bw=Pt(1.5))
tb(s,
   "The gap: no existing paper combines all four layers into a single evaluated Kubernetes system. "
   "Chaudhari (2025) calls for exactly this combination but never builds it. We do.",
   0.35, 6.34, 12.6, 0.46, sz=Pt(12.5), bold=True, col=C_GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — LITERATURE TABLE
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Literature at a Glance", "18 papers — what each contributes and what we take from it")

# Table header
cols_w = [2.35, 2.0, 2.0, 2.0, 3.3]  # Paper | Contribution | Gap | We Use | Has ML/Math
cols_x = [0.18]
for w in cols_w[:-1]:
    cols_x.append(cols_x[-1] + w + 0.03)

headers = ["Paper", "Core Contribution", "Gap Addressed", "What We Include", "Model Type"]
rect(s, 0.18, 1.3, 12.93, 0.42, fill=C_ACCENT)
for hdr, x, w in zip(headers, cols_x, cols_w):
    tb(s, hdr, x+0.04, 1.31, w-0.06, 0.4, sz=Pt(10), bold=True, col=C_BG, align=PP_ALIGN.CENTER)

rows = [
    ("Chaudhari 2025", "K8s fails for AI; proposes 4-component framework", "Gap anchor", "Architecture blueprint", "None — conceptual"),
    ("Jiang Zhi 2025", "Overcommitment study; Clovers simulator", "Eval environment", "Simulator + traces", "None — rule-based sim"),
    ("Coach 2025", "Co-locate jobs with non-overlapping memory peaks", "Overcommit strategy", "Temporal co-location scoring", "Prediction: temporal patterns"),
    ("DRF 2011", "Dominant Resource Fairness — proven fair allocation", "Fairness backbone", "DRF as Score phase", "Fairness algo: LP-proven"),
    ("Res. Central 2017", "RF predicts P95 usage from tenant history", "Prediction model", "RF + tenant features", "Prediction: Random Forest"),
    ("Blagodurov", "cgroups weights: critical vs batch co-location", "SLA enforcement", "Dynamic cgroup controller", "Sched. optimizer: cgroups Eqs"),
    ("Wang & Yang 2025", "LSTM+DQN on K8s; 32.5% util gain", "K8s deploy blueprint", "Plugin architecture", "Prediction (LSTM) + RL (DQN)"),
    ("Doukha 2025", "RF beats LSTM (MAPE 2.65% vs 17.43%)", "Model selection", "RF justified; MAPE metric", "Prediction comparison: RF vs LSTM"),
    ("Priya 2025", "SloPolicy CRD + plugin + Prometheus loop", "K8s architecture", "CRD design; eval metrics", "Sched. optimizer: QoS score fn"),
    ("Kofi 2025", "LSTM on Google Trace; R²=0.99 w/ preprocessing", "Dataset validation", "Trace + preprocessing", "Prediction: LSTM (R²=0.99)"),
    ("Alatawi 2025", "RL MDP serverless; Gini fairness; 98% SLA", "Fairness metric + RL compare", "Gini coefficient metric", "RL policy: MDP allocation"),
    ("Zhao et al. 2021", "Formal admission control + profit optimization", "Admission formalization", "Admission algorithm", "Sched. optimizer: admission LP"),
    ("Liu & Guitart 2025", "In-node DRC; group-aware cgroup; 319% throughput", "In-node enforcement", "DRC post-placement", "Sched. optimizer: in-node cgroup"),
    ("Kovalenko 2024", "Formal discrete K8s optimization model", "Math formalization", "Constraint structure", "Sched. optimizer: discrete LP"),
    ("Atropos 2025", "Throttle culprit, not victim, during overload", "Overload fallback", "Targeted throttling", "None — monitoring-based"),
    ("Perera 2025", "RL model drift + interpretability risks", "Model justification", "RF interpretability arg.", "None — review paper"),
    ("Pinnapareddy 2025", "Bin packing + cost tools in K8s", "Motivation + tooling", "Kubecost for eval", "None — practitioner"),
    ("Patchamatla", "K8s on OpenStack; VM-container deploy", "Deployment context", "Architecture reference", "None — experimental"),
]

row_h = 0.31
y = 1.72
for j, (paper, contrib, gap, use, model) in enumerate(rows):
    bg = C_DARK_BOX if j % 2 == 0 else RGBColor(0x12, 0x22, 0x35)
    rect(s, 0.18, y, 12.93, row_h, fill=bg)
    vals = [paper, contrib, gap, use, model]
    for val, x, w in zip(vals, cols_x, cols_w):
        col = C_ACCENT if x == cols_x[0] else (C_GOLD if x == cols_x[3] else C_LIGHT)
        tb(s, val, x+0.04, y+0.03, w-0.06, row_h-0.04, sz=Pt(8.5), col=col)
    y += row_h

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Methodology", "Simulation-based evaluation on real Google production cluster data")

# Three columns: Dataset, Method, Metrics
for i, (title, col, body) in enumerate([
    ("Dataset", C_TEAL,
     "Google Cluster Trace v3 (2019 release, public)\n\n"
     "▸ Real job submissions from Google's production cluster\n"
     "▸ Contains: job arrival time, declared CPU/memory requests, actual measured usage, job duration, tenant ID, priority class\n"
     "▸ The gap between declared and actual usage is our primary training signal\n"
     "▸ Validated by multiple papers (Kofi 2025, Jiang Zhi 2025)\n\n"
     "Secondary: Azure Public Dataset (Resource Central) for cross-validation"),
    ("Simulation Method", C_GOLD,
     "Discrete-event simulation (no live cluster required)\n\n"
     "▸ Replay Google Cluster Trace jobs through our simulated cluster\n"
     "▸ Simulator components: job arrival queue, node resource tracker, our scheduler decision engine, SLA violation detector\n"
     "▸ Preprocessing: Savitzky-Golay filtering + min-max normalization (Kofi's pipeline)\n"
     "▸ Compare: default K8s scheduler vs. static DRF baseline vs. our full system\n\n"
     "Options: adapt Clovers (Jiang Zhi) or write custom ~300-line Python simulator"),
    ("Evaluation Metrics", C_ACCENT,
     "Primary metrics:\n"
     "▸ Memory utilization: target ≥ 85% (vs ~45–60% baseline)\n"
     "▸ SLA violation rate: target < 5% (Priya 2025 benchmark)\n"
     "▸ Inter-tenant fairness (Gini coefficient): target ~0.10 (vs ~0.25)\n"
     "▸ Prediction accuracy (MAPE): target < 5% (Doukha RF: 2.65%)\n\n"
     "Secondary metrics:\n"
     "▸ P99 latency reduction vs. default K8s\n"
     "▸ Job queue wait time reduction\n"
     "▸ Admission acceptance rate"),
]):
    x = 0.2 + i * 4.35
    rect(s, x, 1.38, 4.1, 0.48, fill=col)
    tb(s, title, x, 1.38, 4.1, 0.48, sz=Pt(14), bold=True, col=C_BG, align=PP_ALIGN.CENTER)
    rect(s, x, 1.86, 4.1, 4.75, fill=C_DARK_BOX, border=col)
    tb(s, body, x+0.12, 1.93, 3.9, 4.6, sz=Pt(10.5))

rect(s, 0.2, 6.75, 12.9, 0.55, fill=RGBColor(0x00, 0x3A, 0x52), border=C_GOLD)
tb(s,
   "No live cloud account needed. Google Cluster Trace v3 is publicly downloadable. "
   "Simulation gives reproducible, measurable results at scale — consistent with published papers in this field (Blagodurov, Coach).",
   0.35, 6.77, 12.6, 0.5, sz=Pt(11.5), col=C_GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EXPECTED OUTCOMES & CONTRIBUTION
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Expected Outcomes & Contribution", "What we deliver and why it matters")

# Left: targets table
rect(s, 0.2, 1.35, 6.5, 0.45, fill=C_ACCENT)
tb(s, "Performance Targets", 0.2, 1.35, 6.5, 0.45, sz=Pt(13), bold=True, col=C_BG, align=PP_ALIGN.CENTER)

target_rows = [
    ("Memory utilization", "≥ 85%", "~45–60%", C_GREEN),
    ("SLA violation rate", "< 5%", "< 80% compliance under load", C_GREEN),
    ("Fairness (Gini)", "~0.10", "~0.25 (unmanaged)", C_GREEN),
    ("Prediction MAPE", "< 5%", "N/A (no prediction)", C_GREEN),
    ("P99 latency reduction", "≥ 40%", "Baseline", C_GOLD),
    ("Job queue wait time", "Reduced", "Unbounded (FCFS)", C_GOLD),
]
y = 1.8
for metric, target, baseline, col in target_rows:
    rect(s, 0.2, y, 6.5, 0.34, fill=RGBColor(0x12, 0x22, 0x35) if target_rows.index((metric, target, baseline, col)) % 2 else C_DARK_BOX)
    tb(s, metric, 0.28, y+0.04, 2.5, 0.28, sz=Pt(10), col=C_LIGHT)
    tb(s, target, 2.78, y+0.04, 1.5, 0.28, sz=Pt(10), bold=True, col=col, align=PP_ALIGN.CENTER)
    tb(s, baseline, 4.28, y+0.04, 2.35, 0.28, sz=Pt(9.5), col=C_GRAY)
    y += 0.34

# Right: novelty + deliverables
rect(s, 7.0, 1.35, 6.1, 0.45, fill=C_GOLD)
tb(s, "What Makes This Novel", 7.0, 1.35, 6.1, 0.45, sz=Pt(13), bold=True, col=C_BG, align=PP_ALIGN.CENTER)
rect(s, 7.0, 1.8, 6.1, 2.28, fill=C_DARK_BOX, border=C_GOLD)
bullets(s, [
    "Predictive admission via RF on tenant history + memory peaks",
    "DRF fairness extended with temporal usage prediction",
    "cgroups-enforced SLA guarantee + Atropos-style fallback",
    "Formal math backbone (Kovalenko) + Gini fairness metric",
    "No existing paper combines all four in one evaluated system",
], 7.12, 1.88, 5.88, 2.1, sz=Pt(11), col=C_LIGHT)

rect(s, 7.0, 4.22, 6.1, 0.45, fill=C_ACCENT)
tb(s, "Project Deliverables", 7.0, 4.22, 6.1, 0.45, sz=Pt(13), bold=True, col=C_BG, align=PP_ALIGN.CENTER)
rect(s, 7.0, 4.67, 6.1, 1.75, fill=C_DARK_BOX, border=C_ACCENT)
bullets(s, [
    "Trained Random Forest prediction model (on Google Cluster Trace)",
    "Custom Kubernetes scheduler plugin (Python / Go prototype)",
    "Discrete-event simulation comparing 3 scheduler strategies",
    "Evaluation report: utilization, SLA, fairness metrics",
], 7.12, 4.75, 5.88, 1.6, sz=Pt(11), col=C_LIGHT)

# Bottom: stakeholder impact
rect(s, 0.2, 6.6, 12.9, 0.72, fill=RGBColor(0x00, 0x2A, 0x3A), border=C_GREEN, bw=Pt(1.5))
tb(s, "Impact", 0.35, 6.62, 2.0, 0.3, sz=Pt(12), bold=True, col=C_GREEN)
tb(s,
   "Cloud providers reduce hardware costs by getting more work from existing servers. "
   "Tenants get fair, predictable resource access and fewer SLA violations. "
   "The scheduling model is platform-agnostic — applicable to any Kubernetes cluster (AWS EKS, GCP GKE, Azure AKS, private clouds).",
   0.35, 6.93, 12.6, 0.35, sz=Pt(11), col=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(r"C:\Users\alric\Documents\Spring2026\Capstone_ClaudeAI\Team_Concept_Slides.pptx")
print("Saved: Team_Concept_Slides.pptx")
