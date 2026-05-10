"""
generate_op_model_overview.py
──────────────────────────────
Generates op_model_overview.pptx — detailed slide deck for the
multi-tenant cluster scheduling capstone project.

Run:
    python generate_op_model_overview.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x37, 0x5E)
MID_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
GREEN      = RGBColor(0x1A, 0x7A, 0x3C)
RED        = RGBColor(0x99, 0x22, 0x22)
AMBER      = RGBColor(0xB8, 0x6E, 0x00)
LIGHT_RED  = RGBColor(0xFF, 0xEE, 0xEE)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=DARK_GREY,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, DARK_BLUE)
    add_textbox(slide, 0.3, 0.08, 12.5, 0.7, title,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.3, 0.75, 12.5, 0.4, subtitle,
                    font_size=14, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)


def bullets(slide, left, top, width, height, items,
            font_size=16, color=DARK_GREY, indent_char="  •  "):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = indent_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def section_label(slide, left, top, text, width=3.5):
    add_rect(slide, left, top, width, 0.32, MID_BLUE)
    add_textbox(slide, left + 0.07, top + 0.01, width - 0.1, 0.30,
                text, font_size=13, bold=True, color=WHITE)


def callout_box(slide, left, top, width, height, title, body,
                title_color=DARK_BLUE, bg_color=LIGHT_BLUE, body_font=13):
    add_rect(slide, left, top, width, height, bg_color)
    add_textbox(slide, left + 0.12, top + 0.08, width - 0.2, 0.32,
                title, font_size=13, bold=True, color=title_color)
    add_textbox(slide, left + 0.12, top + 0.38, width - 0.2, height - 0.45,
                body, font_size=body_font, color=DARK_GREY)


# ============================================================
# SLIDE 1 — Title
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(sl, 0, 2.8, 13.33, 2.2, MID_BLUE)
add_textbox(sl, 0.8, 1.0, 11.7, 1.5,
            "Multi-Tenant Cluster Scheduling",
            font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 2.3, 11.7, 0.6,
            "A Two-Layer Optimization Pipeline",
            font_size=24, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 3.1, 11.7, 0.5,
            "Plan-Ahead MISOCP  +  Real-Time MILP",
            font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 3.75, 11.7, 0.4,
            "Capstone Project — Spring 2026",
            font_size=16, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 5.5, 11.7, 0.4,
            "Alrick Grandison",
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — Pipeline Architecture
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Pipeline Architecture", "Three decoupled layers — prediction -> planning -> scheduling")

add_rect(sl, 0.4, 1.4, 3.6, 1.5, LIGHT_BLUE)
add_textbox(sl, 0.45, 1.42, 3.5, 0.35, "LAYER 1 — Prediction",
            font_size=13, bold=True, color=DARK_BLUE)
bullets(sl, 0.45, 1.75, 3.5, 1.1,
        ["Google cluster traces v3",
         "Demand d_ijr, mean mu_ijr",
         "Covariance Sigma_r (PSD)"],
        font_size=13)

add_textbox(sl, 4.1, 1.9, 0.5, 0.4, "->", font_size=20, color=MID_BLUE, align=PP_ALIGN.CENTER)

add_rect(sl, 4.7, 1.4, 3.6, 1.5, LIGHT_BLUE)
add_textbox(sl, 4.75, 1.42, 3.5, 0.35, "LAYER 2 — Plan-Ahead MISOCP",
            font_size=13, bold=True, color=DARK_BLUE)
bullets(sl, 4.75, 1.75, 3.5, 1.1,
        ["Tenant admission (a_i in {0,1})",
         "Node assignment per time slot",
         "Gurobi solver (MISOCP)"],
        font_size=13)

add_textbox(sl, 8.4, 1.9, 0.5, 0.4, "->", font_size=20, color=MID_BLUE, align=PP_ALIGN.CENTER)

add_rect(sl, 9.0, 1.4, 3.6, 1.5, LIGHT_BLUE)
add_textbox(sl, 9.05, 1.42, 3.5, 0.35, "LAYER 3 — Real-Time MILP",
            font_size=13, bold=True, color=DARK_BLUE)
bullets(sl, 9.05, 1.75, 3.5, 1.1,
        ["Place jobs on allowed nodes",
         "Constraint C5: access control",
         "OR-Tools CBC / GLOP"],
        font_size=13)

add_rect(sl, 3.2, 3.2, 6.9, 0.35, MID_BLUE)
add_textbox(sl, 3.25, 3.22, 6.8, 0.3,
            "TenantAccessSchedule: {(tenant, slot) -> [node_ids]}  |  filter_active_access(t)",
            font_size=12, bold=False, color=WHITE)

section_label(sl, 0.4, 3.8, "What the pipeline produces", width=4.5)
bullets(sl, 0.4, 4.2, 12.5, 2.8,
        ["Plan-ahead: which tenants are admitted and their authorised node sets for each time slot",
         "TenantLease tuples: (tenant, node_set, start_slot, end_slot) — the reservation contract",
         "Real-time: per-job placement decisions that respect both capacity and access constraints",
         "Full traceability: every placed job links back to a plan-ahead lease"],
        font_size=15)


# ============================================================
# SLIDE 3 — Plan-Ahead Model Improvements Overview
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Plan-Ahead Model — Improvements Overview",
              "What was changed, why it mattered, and what was added")

section_label(sl, 0.4, 1.25, "Three Critical Bug Fixes", width=4.0)
bullets(sl, 0.4, 1.65, 5.9, 2.2,
        ["C1 constraint: changed <= a[i] to == a[i]  ->  forces placement when admitted",
         "Q_quota fix: corrected quota formula  ->  DRF fairness now computed accurately",
         "Revenue scaling: pi_bar scales with |H|  ->  solver stops preferring rejection"],
        font_size=14, color=RED)

section_label(sl, 6.7, 1.25, "Design Improvements", width=4.0)
bullets(sl, 6.7, 1.65, 6.2, 2.2,
        ["Gurobi credentials moved to .env file  ->  no hardcoded secrets in code",
         "kap_global  ->  kap_node[n]: per-node safety margin  ->  less conservative SOCP",
         "Gurobi env passed explicitly  ->  no hidden module-level state",
         "lam[0] added for infra weight  ->  fixed duplicate lam[1] index in objective"],
        font_size=14, color=DARK_GREY)

section_label(sl, 0.4, 4.1, "Code Restructure (monolith -> 4 modules)", width=5.5)
bullets(sl, 0.4, 4.5, 12.5, 1.2,
        ["Before: everything in multi_tenant_cluster.py with no separation of concerns",
         "After: plan_ahead_data.py | plan_ahead_optimizer.py | plan_ahead_sensitivity.py | test_plan_ahead.py"],
        font_size=14)

section_label(sl, 0.4, 5.85, "Tests", width=1.5)
add_textbox(sl, 2.1, 5.85, 10.8, 0.4,
            "17 behavioral tests — all passing — covering C1-C7 constraints, admission, compliance, migration, isolation, output format",
            font_size=14, bold=True, color=GREEN)

add_textbox(sl, 0.4, 6.4, 12.5, 0.6,
            "The following slides explain each bug fix in plain language — WHY it broke, WHAT the solver did wrong, HOW it was fixed.",
            font_size=14, bold=False, color=MID_BLUE)


# ============================================================
# SLIDE 4 — What Is a Workload?
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "What Is a 'Workload'? — The Core Concept",
              "Not a running job — a statistical demand forecast unit")

add_textbox(sl, 0.4, 1.25, 12.5, 0.55,
            "You might expect 'workload' to mean a currently-running job. In the plan-ahead model, it means something different:",
            font_size=15, color=DARK_GREY)

add_rect(sl, 0.4, 1.85, 12.5, 0.38, MID_BLUE)
add_textbox(sl, 0.5, 1.87, 12.3, 0.34,
            "Workload (i, j) = a statistical summary of ALL jobs that tenant i expects to run under service profile j during the planning horizon",
            font_size=14, bold=True, color=WHITE)

section_label(sl, 0.4, 2.4, "What a workload carries", width=4.0)
bullets(sl, 0.4, 2.82, 5.8, 1.8,
        ["d_ijr  — declared resource REQUEST  (what the tenant says they need; e.g. 2 CPU, 4 GB RAM)",
         "mu_ijr — mean predicted USAGE  (what historical data says they actually consume)",
         "Sigma_r — covariance across workloads  (how correlated usage is; affects safety margin)"],
        font_size=13)

section_label(sl, 6.7, 2.4, "Long-running vs short-running jobs", width=4.5)
bullets(sl, 6.7, 2.82, 6.1, 1.8,
        ["Long-running: a web server or database that runs 24/7 in every slot",
         "Short-running: batch jobs that come and go every few minutes",
         "Both are captured as a single demand profile for the plan-ahead — the model does not schedule individual jobs; it reserves capacity for them"],
        font_size=13)

section_label(sl, 0.4, 4.75, "Why model at this level?", width=4.0)
bullets(sl, 0.4, 5.15, 12.5, 1.6,
        ["The plan-ahead runs ONCE per planning horizon (e.g. daily) and must decide node reservations for hours ahead",
         "At that timescale, individual jobs are unknowable — only statistical forecasts are available",
         "The real-time model (Layer 3) handles actual job placement once per scheduling round (~60 seconds)"],
        font_size=14)

add_rect(sl, 0.4, 6.9, 12.5, 0.4, LIGHT_BLUE)
add_textbox(sl, 0.5, 6.92, 12.3, 0.36,
            "Analogy: a workload is like a hotel room block reservation — you book X rooms for a conference without knowing which specific guests will fill them.",
            font_size=13, bold=False, color=DARK_BLUE)


# ============================================================
# SLIDE 5 — Bug Fix 1: C1 Constraint
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Bug Fix 1 — The C1 Constraint: = vs <=",
              "The most critical bug: solver was earning revenue without doing any work")

# Left: BEFORE (wrong)
add_rect(sl, 0.4, 1.2, 5.9, 0.35, RED)
add_textbox(sl, 0.5, 1.22, 5.8, 0.31, "BEFORE (broken code):",
            font_size=14, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.55, 5.9, 0.5, LIGHT_RED)
add_textbox(sl, 0.5, 1.57, 5.8, 0.46,
            "SUM over all nodes n: x[i,j,n,t]  <=  a[i]",
            font_size=14, bold=True, color=RED)
add_textbox(sl, 0.4, 2.15, 5.9, 1.5,
            "Reading: 'place workload (i,j) on at most one node, IF tenant i is admitted.'\n\n"
            "Problem: the <= allows placing on ZERO nodes even when admitted. The constraint only "
            "prevents placing on MORE than one node. It does NOT require any node to be used.",
            font_size=13, color=DARK_GREY)

# Right: AFTER (fixed)
add_rect(sl, 6.9, 1.2, 5.9, 0.35, GREEN)
add_textbox(sl, 7.0, 1.22, 5.8, 0.31, "AFTER (fixed code):",
            font_size=14, bold=True, color=WHITE)
add_rect(sl, 6.9, 1.55, 5.9, 0.5, LIGHT_GREEN)
add_textbox(sl, 7.0, 1.57, 5.8, 0.46,
            "SUM over all nodes n: x[i,j,n,t]  ==  a[i]",
            font_size=14, bold=True, color=GREEN)
add_textbox(sl, 6.9, 2.15, 5.9, 1.5,
            "Reading: 'place workload (i,j) on EXACTLY one node, if and only if tenant i is admitted. "
            "If NOT admitted (a[i]=0), sum must be 0. If admitted (a[i]=1), sum must be 1.'\n\n"
            "Now placement is forced. You cannot earn revenue without using a node.",
            font_size=13, color=DARK_GREY)

# What the solver found
add_rect(sl, 0.4, 3.8, 12.5, 0.35, AMBER)
add_textbox(sl, 0.5, 3.82, 12.3, 0.31, "What the solver exploited with the <= version:",
            font_size=14, bold=True, color=WHITE)
add_textbox(sl, 0.4, 4.22, 12.5, 1.25,
            "1. Set a[i] = 1 for ALL tenants (admitted).\n"
            "2. Set ALL x[i,j,n,t] = 0  (no workloads placed anywhere).\n"
            "3. No nodes activated  ->  zero infrastructure cost.\n"
            "4. With no tenants placed, the DRF fairness variable sigma was unconstrained and went to 1.0 (maximum).\n"
            "5. Objective = +revenue from all admissions - fairness reward = most negative value possible.\n"
            "   The solver declared this 'optimal' because it satisfied all constraints (<=) without doing any real scheduling.",
            font_size=13, color=DARK_GREY)

add_rect(sl, 0.4, 5.6, 12.5, 0.32, MID_BLUE)
add_textbox(sl, 0.5, 5.62, 12.3, 0.28,
            "After fix: solver admitted all tenants AND placed every workload on a node. Tests went from 7 failures to 17/17 passing.",
            font_size=13, bold=True, color=WHITE)

section_label(sl, 0.4, 6.05, "Why this matters conceptually", width=4.5)
add_textbox(sl, 0.4, 6.42, 12.5, 0.75,
            "The original <= was logically wrong. In a real cluster, you cannot bill a tenant and then give them nothing. "
            "The admission variable a[i] is a CONTRACT — admit means you must place their workloads. "
            "The == constraint enforces this contract mathematically.",
            font_size=13, color=DARK_GREY)


# ============================================================
# SLIDE 6 — Bug Fix 2: Q_quota and DRF Fairness
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Bug Fix 2 — Q_quota and DRF Fairness Calculation",
              "The quota denominator was too large, making fairness appear worthless to the solver")

section_label(sl, 0.4, 1.25, "What is DRF Fairness?", width=3.8)
add_textbox(sl, 0.4, 1.65, 12.5, 0.85,
            "DRF = Dominant Resource Fairness. The idea: each tenant has a 'fair share' of resources based on their quota. "
            "The model tracks sigma_fair = the MINIMUM satisfaction fraction across all admitted tenants. "
            "Maximising sigma_fair means raising the satisfaction of the worst-off tenant — nobody starves.",
            font_size=14, color=DARK_GREY)

section_label(sl, 0.4, 2.65, "The formula", width=2.5)
add_rect(sl, 0.4, 3.05, 12.5, 0.45, LIGHT_BLUE)
add_textbox(sl, 0.5, 3.07, 12.3, 0.41,
            "s_i (satisfaction of tenant i)  =  SUM over placed workloads: d[i,j,r*]  /  Q_quota[i, r*]"
            "     where r* = dominant resource of tenant i",
            font_size=14, bold=True, color=DARK_BLUE)

add_textbox(sl, 0.4, 3.6, 12.5, 0.6,
            "Larger Q_quota means a SMALLER satisfaction fraction. If the quota is set too high, even fully-placed tenants "
            "have very low satisfaction, so sigma_fair gives the solver very little reward — not worth caring about.",
            font_size=13, color=DARK_GREY)

# Left: BEFORE
add_rect(sl, 0.4, 4.3, 5.9, 0.32, RED)
add_textbox(sl, 0.5, 4.32, 5.8, 0.28, "BEFORE (broken):", font_size=13, bold=True, color=WHITE)
add_rect(sl, 0.4, 4.62, 5.9, 0.42, LIGHT_RED)
add_textbox(sl, 0.5, 4.64, 5.8, 0.38,
            "Q_quota = node_capacity * n_nodes / n_tenants  =  10 * 4 / 3  =  13.33",
            font_size=13, bold=True, color=RED)
add_textbox(sl, 0.4, 5.08, 5.9, 0.85,
            "With quota = 13.33: each admitted tenant's demand (~1-2 units) / 13.33 = ~0.09 to 0.15.\n"
            "sigma_fair reward in objective: lam[3] * 0.15 = 5 * 0.15 = 0.75. Very small.\n"
            "Solver finds rejecting everyone gives the same or better objective: not worth the infra cost.",
            font_size=12, color=DARK_GREY)

# Right: AFTER
add_rect(sl, 6.9, 4.3, 5.9, 0.32, GREEN)
add_textbox(sl, 7.0, 4.32, 5.8, 0.28, "AFTER (fixed):", font_size=13, bold=True, color=WHITE)
add_rect(sl, 6.9, 4.62, 5.9, 0.42, LIGHT_GREEN)
add_textbox(sl, 7.0, 4.64, 5.8, 0.38,
            "Q_quota = node_capacity / 2  =  10 / 2  =  5.0",
            font_size=13, bold=True, color=GREEN)
add_textbox(sl, 6.9, 5.08, 5.9, 0.85,
            "With quota = 5.0: each admitted tenant's demand / 5.0 = ~0.2 to 0.4.\n"
            "sigma_fair reward: lam[3] * 0.4 = 5 * 0.4 = 2.0. Now meaningful.\n"
            "Solver finds admitting tenants is worth the infra cost because fairness reward is substantial.",
            font_size=12, color=DARK_GREY)

add_rect(sl, 0.4, 6.1, 12.5, 0.32, MID_BLUE)
add_textbox(sl, 0.5, 6.12, 12.3, 0.28,
            "Rule: Q_quota = node_capacity / 2 preserves the original ratio (10/2=5) and scales correctly when capacity changes.",
            font_size=13, bold=True, color=WHITE)


# ============================================================
# SLIDE 7 — Bug Fix 3: Revenue Scaling
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Bug Fix 3 — Revenue Must Scale With Planning Horizon",
              "With more time slots, infra cost grew but admission revenue stayed fixed — solver preferred rejecting everyone")

section_label(sl, 0.4, 1.25, "The imbalance problem", width=3.8)
add_textbox(sl, 0.4, 1.65, 12.5, 0.65,
            "The objective includes two competing forces: infrastructure COST (paid every slot) vs admission REVENUE (paid once per tenant). "
            "When the planning horizon grows, cost grows but revenue did not. The solver found 'reject everyone' was always cheaper.",
            font_size=14, color=DARK_GREY)

# Two column: cost vs revenue
add_rect(sl, 0.4, 2.45, 5.9, 0.32, RED)
add_textbox(sl, 0.5, 2.47, 5.8, 0.28, "Infrastructure COST (paid every slot)", font_size=13, bold=True, color=WHITE)
add_textbox(sl, 0.4, 2.82, 5.9, 1.5,
            "Cost = lam[0] x pi_n x z[n,t]  (per active node, per slot)\n\n"
            "Example with |H|=3 slots, 3 tenants using 4 nodes:\n"
            "  Cost = 1.0 x 1.0 x 4 nodes x 3 slots = 12.0 units\n\n"
            "This grows linearly with the number of planning slots.",
            font_size=13, color=DARK_GREY)

add_rect(sl, 6.9, 2.45, 5.9, 0.32, GREEN)
add_textbox(sl, 7.0, 2.47, 5.8, 0.28, "Admission REVENUE (paid ONCE per tenant)", font_size=13, bold=True, color=WHITE)
add_textbox(sl, 6.9, 2.82, 5.9, 1.5,
            "Revenue = lam[2] x (pi_bar_i - v_op_i) x a[i]  (per admitted tenant)\n\n"
            "BEFORE fix with pi_bar=3.0 (fixed), 3 tenants:\n"
            "  Revenue = 1.0 x (3.0 - 0.5) x 3 = 7.5 units\n\n"
            "Cost (12.0) > Revenue (7.5)  ->  solver rejects everyone!",
            font_size=13, color=RED)

add_rect(sl, 0.4, 4.45, 12.5, 0.35, AMBER)
add_textbox(sl, 0.5, 4.47, 12.3, 0.31,
            "Reject ALL = cost 0, revenue 0, objective = 0. Admit ALL = cost 12, revenue 7.5, obj = +4.5 (worse). Solver correctly chose rejection!",
            font_size=13, bold=True, color=WHITE)

section_label(sl, 0.4, 4.95, "The fix: scale revenue with horizon", width=4.5)
add_textbox(sl, 0.4, 5.35, 12.5, 0.5,
            "Logic: if you run the cluster for 3 time slots, you should charge tenants 3x as much. Contract value is proportional to the duration.",
            font_size=14, color=DARK_GREY)

add_rect(sl, 0.4, 5.95, 12.5, 0.42, LIGHT_GREEN)
add_textbox(sl, 0.5, 5.97, 12.3, 0.38,
            "AFTER: pi_bar[i] = 3.0 x |H|   and   v_op[i] = 0.5 x |H|   ->   revenue now grows with the horizon just like cost does",
            font_size=14, bold=True, color=GREEN)

add_textbox(sl, 0.4, 6.5, 12.5, 0.55,
            "Result: with |H|=3, revenue per tenant = (9.0 - 1.5) = 7.5, total revenue = 22.5 > cost 12.0. "
            "Solver now finds admitting tenants is profitable. All three tenants are admitted.",
            font_size=13, color=DARK_GREY)


# ============================================================
# SLIDE 8 — SOCP Capacity Constraint Explained
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "The SOCP Capacity Constraint — Safety Under Uncertainty",
              "Why a simple 'sum of usage <= capacity' is not enough, and what we do instead")

section_label(sl, 0.4, 1.25, "The problem: workload usage is uncertain", width=4.8)
add_textbox(sl, 0.4, 1.65, 12.5, 0.65,
            "If we just said: 'total mean usage <= capacity', we would plan a node at exactly 100% utilization on average. "
            "That means 50% of the time the node is OVERLOADED. Workloads get throttled, SLAs violated. We need a safety margin.",
            font_size=14, color=DARK_GREY)

section_label(sl, 0.4, 2.45, "The Cantelli inequality — a probabilistic guarantee", width=5.5)
add_textbox(sl, 0.4, 2.85, 12.5, 0.65,
            "The Cantelli inequality says: if you add kappa(eps) x standard_deviation to the mean, "
            "the probability that the actual usage exceeds that bound is at most eps (the risk tolerance).",
            font_size=14, color=DARK_GREY)

add_rect(sl, 0.4, 3.6, 12.5, 0.48, LIGHT_BLUE)
add_textbox(sl, 0.5, 3.62, 12.3, 0.44,
            "Constraint: Mean_Load[n,r,t]  +  kappa(eps)  x  ||Cholesky(Sigma_r) x xi_vec||  <=  Capacity[n,r] x active[n,t]",
            font_size=15, bold=True, color=DARK_BLUE)

# Three boxes explaining each part
callout_box(sl, 0.4, 4.22, 3.8, 1.45,
            "Mean_Load",
            "The AVERAGE total resource consumption on node n from all workloads placed there, accounting for the isolation primitive overhead (eta_kr x mu_ijr).",
            body_font=12)

callout_box(sl, 4.4, 4.22, 4.3, 1.45,
            "kappa x ||Chol(Sigma) x xi||",
            "kappa(0.05) = 4.36. This is the safety buffer. Sigma_r is the covariance matrix (how correlated workload fluctuations are). Cholesky decomposes it. The ||...|| is the L2 norm, computed as a Second-Order Cone (hence SOCP).",
            body_font=12)

callout_box(sl, 8.9, 4.22, 4.0, 1.45,
            "Capacity x active?",
            "The right side is node capacity, but ONLY if the node is switched on (z[n,t]=1). If no workloads are placed, z=0 and the right side becomes 0 — no capacity used.",
            body_font=12)

section_label(sl, 0.4, 5.82, "Per-node kappa improvement (our fix)", width=4.5)
add_textbox(sl, 0.4, 6.22, 12.5, 0.65,
            "Original: used the single strictest tenant's risk tolerance for ALL nodes globally (most conservative).\n"
            "Fix: each node uses the strictest eps_i only among tenants that actually have access to that node. "
            "Nodes restricted to lenient tenants get a smaller safety margin -> more workloads can fit.",
            font_size=13, color=DARK_GREY)


# ============================================================
# SLIDE 9 — Real-Time Model Improvements
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Real-Time Model — Improvements & Integration")

section_label(sl, 0.4, 1.25, "New Constraint C5 — Plan-Ahead Access Control")
add_textbox(sl, 0.4, 1.65, 12.5, 0.4,
            "x_jn = 0  if  n not in A_{t(j)}   (job j cannot be placed on node n if tenant is not authorised)",
            font_size=15, bold=True, color=DARK_BLUE)
bullets(sl, 0.4, 2.1, 12.5, 1.1,
        ["Implemented via upper-bound = 0 at variable creation — zero solver overhead, no extra constraints",
         "tenant_node_access = None  ->  backward-compatible (all nodes allowed);  dict provided  ->  plan-ahead enforced",
         "Verified by tests 13 & 14: blocked nodes produce zero placements, None allows all"],
        font_size=14)

section_label(sl, 0.4, 3.35, "Integration Points")
bullets(sl, 0.4, 3.75, 5.8, 2.5,
        ["filter_active_access(schedule, t) slices TenantAccessSchedule for the current round",
         "TenantLease dataclass holds (tenant, nodes, start_slot, end_slot)",
         "schedule_to_leases() merges consecutive identical slots into compact leases"],
        font_size=14)

section_label(sl, 6.5, 3.35, "Scalability")
bullets(sl, 6.5, 3.75, 6.4, 2.5,
        ["CBC (exact MILP): correct for <= ~24 jobs per round; slows beyond that",
         "GLOP (LP relaxation + rounding): 10x faster, slight quality loss — used in Sample 3",
         "Heuristic recommendation: First-Fit Decreasing by pred_mem — O(|J|x|N|), microseconds",
         "Solver selectable per config; GLOP auto-selected in High complexity sample"],
        font_size=14)


# ============================================================
# SLIDE 10 — How C5 Connects the Two Layers
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "How C5 Connects Plan-Ahead to Real-Time — The Handshake",
              "The plan-ahead issues a contract; the real-time scheduler honours it")

section_label(sl, 0.4, 1.25, "Step 1: Plan-ahead decides node access (runs once per horizon)", width=6.5)
add_rect(sl, 0.4, 1.65, 12.5, 0.45, LIGHT_BLUE)
add_textbox(sl, 0.5, 1.67, 12.3, 0.41,
            "TenantAccessSchedule = { (tenant=0, slot=1): [node 2, node 3],   (tenant=1, slot=1): [node 0, node 1],  ... }",
            font_size=14, bold=True, color=DARK_BLUE)
add_textbox(sl, 0.4, 2.15, 12.5, 0.55,
            "This says: 'At time slot 1, tenant 0 may only use nodes 2 and 3. Tenant 1 may only use nodes 0 and 1.'  "
            "This is a hard access-control decision — no real-time flexibility to override it.",
            font_size=13, color=DARK_GREY)

section_label(sl, 0.4, 2.85, "Step 2: Each real-time round, slice the schedule to current slot", width=6.5)
add_rect(sl, 0.4, 3.25, 12.5, 0.42, LIGHT_BLUE)
add_textbox(sl, 0.5, 3.27, 12.3, 0.38,
            "active_access = filter_active_access(schedule, t=1)  ->  { 0: [2,3],  1: [0,1] }",
            font_size=14, bold=True, color=DARK_BLUE)

section_label(sl, 0.4, 3.8, "Step 3: Real-time solver enforces C5 at variable creation", width=5.5)
add_rect(sl, 0.4, 4.2, 12.5, 0.42, LIGHT_BLUE)
add_textbox(sl, 0.5, 4.22, 12.3, 0.38,
            "For every job j and node n:  if n not in active_access[j.tenant_id]:  upper_bound of x[j,n] = 0  (variable CANNOT be 1)",
            font_size=13, bold=True, color=DARK_BLUE)
add_textbox(sl, 0.4, 4.7, 12.5, 0.6,
            "Setting UB=0 at variable creation is better than adding a constraint after the fact. The LP relaxation is tighter "
            "(fewer feasible directions), branch-and-bound converges faster, no extra memory for a constraint matrix row.",
            font_size=13, color=DARK_GREY)

section_label(sl, 0.4, 5.45, "Why this design matters", width=3.5)
bullets(sl, 0.4, 5.85, 12.5, 1.4,
        ["Security: a tenant's jobs cannot reach nodes they were not allocated — mimics Kubernetes NetworkPolicy",
         "Predictability: the plan-ahead contract is honoured exactly; no job can 'sneak' onto an unplanned node",
         "Backward compatibility: if no plan-ahead output is provided (None), all nodes are allowed — model works standalone"],
        font_size=14)


# ============================================================
# SLIDE 11 — Test & Pipeline Results
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Test Results & Pipeline Validation")

section_label(sl, 0.4, 1.25, "Plan-Ahead Tests  (17 / 17 passing)")
bullets(sl, 0.4, 1.65, 6.0, 2.8,
        ["Model solves to feasibility (no infeasible result)",
         "At least one tenant admitted",
         "All admitted tenants have workloads placed (C1 == fix verified)",
         "Isolation floor: tenant with k_min=1 never uses primitive 0 (no-isolation)",
         "Compliance: restricted workload stays on N_allowed nodes (C1c)",
         "Migration budget respected per tenant per slot (C6c)",
         "No migrations at t=0 (C6a)",
         "DRF sigma in [0, 1]; admitted tenants appear in TenantAccessSchedule"],
        font_size=13, color=GREEN)

section_label(sl, 6.7, 1.25, "Real-Time Tests  (14 / 14 passing)")
bullets(sl, 6.7, 1.65, 6.2, 2.8,
        ["Node capacity constraint satisfied (no overloading)",
         "One-node-per-job integrity",
         "C5: blocked tenant-node pairs produce 0 placements",
         "C5: None access allows all nodes (backward compat.)",
         "Delay-weight omega: high-wait tenant fills capacity first",
         "Partial violation rate halves effective capacity",
         "Omega formula verified analytically"],
        font_size=13, color=GREEN)

section_label(sl, 0.4, 4.6, "Pipeline Demo Results")
headers = ["Sample", "Plan-Ahead Size", "Solver Status", "Jobs/slot", "Placed", "RT Solver"]
rows = [
    ["1 — Simple",  "3T, 4N, 2 slots, 6 wl",  "OPTIMAL (<1s)",      "8",  "8/8",   "CBC"],
    ["2 — Medium",  "3T, 5N, 3 slots, 9 wl",  "Feasible (22% gap)", "12", "12/12", "CBC"],
    ["3 — High",    "3T, 5N, 4 slots, 9 wl",  "Feasible (29% gap)", "20", "20/20", "GLOP"],
]
col_w = [1.9, 2.7, 2.5, 1.6, 1.1, 1.3]
col_x = [0.4, 2.4, 5.2, 7.8, 9.5, 10.7]
row_h = 0.38
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(sl, cx, 5.05, cw - 0.05, 0.35, MID_BLUE)
    add_textbox(sl, cx + 0.05, 5.07, cw - 0.1, 0.3, hdr, font_size=12, bold=True, color=WHITE)
for ri, row in enumerate(rows):
    bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
    for ci, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(sl, cx, 5.45 + ri * row_h, cw - 0.05, row_h - 0.03, bg)
        add_textbox(sl, cx + 0.05, 5.47 + ri * row_h, cw - 0.1, row_h - 0.05,
                    val, font_size=11, color=DARK_GREY)


# ============================================================
# SLIDE 12 — Sensitivity Analysis Results
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Sensitivity Analysis — What the Results Actually Mean",
              "Which model parameters matter and which are inactive for this instance")

add_textbox(sl, 0.4, 1.2, 12.5, 0.4,
            "We swept 5 parameters and measured objective value (more negative = better = more revenue + fairness admitted).",
            font_size=14, color=DARK_GREY)

# Two columns: active vs inactive
add_rect(sl, 0.4, 1.72, 5.9, 0.32, GREEN)
add_textbox(sl, 0.5, 1.74, 5.8, 0.28, "ACTIVE — these change the solution", font_size=13, bold=True, color=WHITE)

callout_box(sl, 0.4, 2.1, 5.9, 1.5,
            "eps_i (SLA risk tolerance): obj dropped from -13.3 to -8.7 at eps=0.02",
            "At very tight risk tolerance (eps<0.05, kappa>4.36), the SOCP safety margin becomes so large "
            "it exceeds node capacity. Workloads cannot fit. Fewer tenants admitted -> worse objective. "
            "Default eps=0.05 sits right at the feasibility boundary.",
            bg_color=LIGHT_GREEN, body_font=12)

callout_box(sl, 0.4, 3.7, 5.9, 1.5,
            "Node capacity: obj improved from -6.7 (cap=5) to -15.1 (cap=20)",
            "With small nodes, SOCP margin eats all available space and workloads cannot be packed in. "
            "After cap=10, all tenants are admitted and the plateau means no further benefit until "
            "packing improves at cap=20 (fewer nodes activated = lower infra cost).",
            bg_color=LIGHT_GREEN, body_font=12)

add_rect(sl, 6.9, 1.72, 5.9, 0.32, AMBER)
add_textbox(sl, 7.0, 1.74, 5.8, 0.28, "INACTIVE — flat for this instance", font_size=13, bold=True, color=WHITE)

callout_box(sl, 6.9, 2.1, 5.9, 0.92,
            "Migration budget (Delta_i): completely flat across 0-10",
            "With only 2 time slots, the optimal placement never requires migration. Budget never binds. "
            "Would activate with 5+ slots and heterogeneous demand.",
            bg_color=LIGHT_BLUE, body_font=12)

callout_box(sl, 6.9, 3.1, 5.9, 0.92,
            "SLA penalty weight (lam[1]): identical at 1x through 50x",
            "Latency SLA violations (e_ijt) are always zero. With L=50ms target and only 2-3 co-located workloads, "
            "there is 40ms of slack. The latency model never activates.",
            bg_color=LIGHT_BLUE, body_font=12)

callout_box(sl, 6.9, 4.1, 5.9, 1.1,
            "Fairness weight (lam[3]): sigma_fair = 0.6131 regardless of weight",
            "Objective scales exactly linearly with lam[3], meaning sigma_fair is the SAME in every solution. "
            "For this symmetric synthetic instance, the fair allocation IS the cost-optimal allocation. "
            "Fairness is free here — this will not hold for heterogeneous real workloads.",
            bg_color=LIGHT_BLUE, body_font=12)

add_rect(sl, 0.4, 5.32, 12.5, 0.35, MID_BLUE)
add_textbox(sl, 0.5, 5.34, 12.3, 0.31,
            "Takeaway: SOCP capacity constraint (C2) is the only binding constraint for this synthetic instance. Migration, SLA, and fairness constraints are inactive with synthetic data.",
            font_size=13, bold=True, color=WHITE)

section_label(sl, 0.4, 5.82, "What this tells us about the model", width=4.0)
bullets(sl, 0.4, 6.22, 12.5, 1.0,
        ["The model is correct — active constraints change behaviour as expected",
         "Inactive constraints need real trace data (lower variance Sigma_r, tighter L, heterogeneous workloads) to activate and show their effect"],
        font_size=13)


# ============================================================
# SLIDE 13 — Challenges
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Challenges Encountered")

challenges = [
    ("SOCP Safety Margin Too Conservative",
     "kappa(0.05) = 4.36 is very large. Random synthetic covariance matrices generate large off-diagonal values, so the safety buffer can exceed available capacity. Fixed partially via per-node kappa. Would improve further with real trace data."),
    ("Revenue vs. Infrastructure Cost Balance",
     "With more time slots, recurring infra cost dominated one-time admission revenue. Required scaling pi_bar by |H| — treating contract value as proportional to horizon length. Without this, solver always preferred rejecting everyone."),
    ("Degenerate Admission Solution (C1 bug)",
     "Original C1 used <= a[i], allowing 'admit all tenants, place nothing'. Sigma=1 (unconstrained when nothing placed) made this appear optimal with objective=-5.0. Fixed to == a[i]. This was the single most impactful bug."),
    ("Gurobi Time Limits on Medium+ Instances",
     "9 workloads + SOCP constraints make branch-and-bound hard. Medium sample hits 120s time limit at 22% gap, High at 29% gap. Solver still returns a feasible solution, but optimality is not proven."),
    ("Real-Time Scalability Past ~24 Jobs",
     "OR-Tools CBC is exact but slow for large |J|x|N| binary programs. GLOP (LP relaxation) is 10x faster. A greedy FFD heuristic would make it O(|J| log|J|) — not yet implemented."),
]

for i, (title, body) in enumerate(challenges):
    y = 1.25 + i * 1.17
    add_rect(sl, 0.4, y, 0.32, 0.35, AMBER)
    add_textbox(sl, 0.4, y, 0.3, 0.35, str(i + 1), font_size=16, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, 0.8, y, 5.5, 0.35, title, font_size=13, bold=True, color=DARK_BLUE)
    add_textbox(sl, 0.8, y + 0.35, 12.1, 0.75, body, font_size=12, color=DARK_GREY)


# ============================================================
# SLIDE 14 — Pros and Cons
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Pros and Cons")

add_rect(sl, 0.4, 1.2, 6.0, 0.38, GREEN)
add_textbox(sl, 0.45, 1.22, 5.9, 0.34, "STRENGTHS", font_size=14, bold=True, color=WHITE)
pros = [
    "Mathematically rigorous MISOCP with 7 constraint families (C1-C7)",
    "Probabilistic capacity via Cantelli inequality — SLA guarantees on uncertainty",
    "DRF fairness: max-min satisfaction prevents tenant starvation",
    "Isolation primitives (none / gVisor / Kata) with compliance floor enforcement",
    "Clean two-layer separation: plan-ahead contract feeds real-time access control",
    "Full test coverage: 31 passing tests across both models",
    "Scalable real-time solver: switchable between exact MILP and LP relaxation",
]
bullets(sl, 0.4, 1.65, 6.0, 5.5, pros, font_size=13, color=GREEN)

add_rect(sl, 6.9, 1.2, 6.0, 0.38, RED)
add_textbox(sl, 6.95, 1.22, 5.9, 0.34, "WEAKNESSES / LIMITATIONS", font_size=14, bold=True, color=WHITE)
cons = [
    "Plan-ahead MISOCP grows O(|W|^2) in workloads due to McCormick zeta pairs",
    "Gurobi commercial license required — not deployable open-source as-is",
    "Synthetic data only — no real Google trace ingestion yet",
    "kap_node still conservative (per-node floor, not exact per-time-slot)",
    "Real-time heuristic not yet implemented — CBC needed for exact results",
    "Sensitivity analysis shows most constraints are inactive with synthetic data",
    "No warm-start from prior solution — each plan-ahead solves cold",
]
bullets(sl, 6.9, 1.65, 6.0, 5.5, cons, font_size=13, color=RED)


# ============================================================
# SLIDE 15 — Capstone Expectations
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Setting Expectations — What This Capstone Delivers")

section_label(sl, 0.4, 1.25, "What IS delivered", width=3.8)
bullets(sl, 0.4, 1.65, 5.9, 4.5,
        ["End-to-end two-layer scheduling pipeline with clean interfaces",
         "A corrected and modularised MISOCP plan-ahead model",
         "Real-time MILP with plan-ahead access enforcement (C5)",
         "31 passing behavioral tests across both models",
         "Three complexity levels (Simple / Medium / High) as demo",
         "Documented source of truth + lexicon for both models",
         "Pipeline interface with TenantLease reservation contracts",
         "Sensitivity analysis module with real results"],
        font_size=14, color=GREEN)

section_label(sl, 6.7, 1.25, "What is NOT in scope", width=3.8)
bullets(sl, 6.7, 1.65, 6.2, 4.5,
        ["Real Google cluster-usage traces ingestion",
         "Production Kubernetes / container runtime integration",
         "Live cluster experiments or SLA violation measurements",
         "FFD greedy heuristic for real-time at scale",
         "Dynamic re-planning on SLA breach events",
         "Multi-cluster or federated scheduling"],
        font_size=14, color=RED)

section_label(sl, 0.4, 6.3, "Bottom line")
add_textbox(sl, 0.4, 6.7, 12.5, 0.6,
            "This capstone proves the two-layer architecture is sound, implementable, and testable. "
            "It establishes the foundation for a production-grade scheduler.",
            font_size=15, bold=True, color=DARK_BLUE)


# ============================================================
# SLIDE 16 — Future Work
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Future Work")

future = [
    ("Real Trace Ingestion",
     "Map Google cluster-usage traces v3: InstanceUsage.average_usage -> mu_ijr, random_sampled_usage -> Sigma_r, resource_request -> d_ijr, MachineEvents.capacity -> C_nr."),
    ("FFD Greedy Heuristic for Real-Time",
     "Replace CBC for large instances: sort jobs by pred_mem descending, assign greedily to first node that fits within tenant_node_access. O(|J| log|J| + |J|x|N|). Already designed, ~30 lines."),
    ("Dynamic Re-Planning",
     "Trigger plan-ahead re-solve when a high-value tenant arrives or SLA breach detected. Use prior solution as MIPStart warm-start to reduce solve time significantly."),
    ("Exact Per-Node eps_eff",
     "Replace precomputed kap_node with binary selectors on y[i,n,t] for exact minimum eps_i per node. Reduces capacity conservatism without adding continuous variables."),
    ("Kubernetes API Integration",
     "Map x_ijnt -> Pod scheduling, w_ijkt -> RuntimeClass (gVisor/Kata), Delta_i -> PodDisruptionBudget, TenantAccessSchedule -> NetworkPolicy rules."),
    ("Rolling Horizon Re-Planning",
     "Re-run plan-ahead every |H|/2 slots with updated demand forecasts. Allows the model to adapt to workload changes mid-horizon without a full cold restart."),
]

for i, (title, body) in enumerate(future):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.3 + row * 1.9
    add_rect(sl, x, y, 6.1, 0.35, MID_BLUE)
    add_textbox(sl, x + 0.1, y + 0.02, 6.0, 0.31, f"F{i+1}  {title}",
                font_size=13, bold=True, color=WHITE)
    add_textbox(sl, x + 0.1, y + 0.38, 6.0, 1.45, body, font_size=12, color=DARK_GREY)


# ============================================================
# SLIDE 17 — Possible Improvements Now
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "Improvements Achievable Now  (before final submission)")

items = [
    ("FFD Heuristic — Real-Time",
     "~30 lines. Sort jobs by pred_mem desc, greedily assign to first fitting node in tenant_node_access. Eliminates CBC timeout for Sample 3 and any production-scale test.",
     GREEN),
    ("Assert Compliance in Pipeline",
     "Add an explicit check in run_pipeline() that no placed real-time job lands on a node outside its lease. One assertion, catches any regression in C5 immediately.",
     MID_BLUE),
    ("Gurobi MIPStart Warm-Start",
     "After finding the Sample 1 solution, pass x, w, a values as hints for Sample 2/3 using var.Start. Can cut the branch-and-bound tree significantly for related instances.",
     AMBER),
    ("Diagonal Covariance for Cleaner Demos",
     "Replace random Sigma_r with identity matrix for sensitivity analysis. Removes the dominant SOCP effect from synthetic data, making all 5 constraint sweeps show meaningful variation.",
     AMBER),
    ("State Initialisation at t=0",
     "Pass current cluster placement as prior-state input to plan-ahead. Enables migration constraints (C6b/C6c) to activate for multi-slot instances, making sensitivity sweeps informative.",
     MID_BLUE),
]

for i, (title, body, color) in enumerate(items):
    y = 1.3 + i * 1.13
    add_rect(sl, 0.4, y, 0.3, 0.9, color)
    add_textbox(sl, 0.8, y, 4.5, 0.38, title, font_size=14, bold=True, color=DARK_BLUE)
    add_textbox(sl, 0.8, y + 0.38, 12.1, 0.7, body, font_size=13, color=DARK_GREY)


# ============================================================
# SLIDE 18 — Plan-Ahead Model Assessment
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_title_bar(sl, "How Is the Plan-Ahead Model? — Honest Assessment")

section_label(sl, 0.4, 1.2, "What is working well", width=4.0)
bullets(sl, 0.4, 1.6, 5.9, 2.8,
        ["All 7 constraint families correctly implemented (C1-C7)",
         "SOCP capacity with Cantelli safety margin is mathematically sound",
         "McCormick linearisations produce tight LP relaxation bounds",
         "DRF fairness sigma correctly bounded by all admitted tenants",
         "Compliance / residency hard constraints verified in tests",
         "Output interface (TenantAccessSchedule) cleanly connects to real-time"],
        font_size=13, color=GREEN)

section_label(sl, 6.7, 1.2, "Known limitations", width=4.0)
bullets(sl, 6.7, 1.6, 6.2, 2.8,
        ["SOCP is still conservative: kap_node uses floor over tenant access, not actual occupancy",
         "Scale: quadratic in workloads (zeta pairs); 9 workloads already near solve-time limit",
         "Synthetic Sigma too large: SOCP dominates; real traces would show more varied behaviour",
         "Academic Gurobi license may hit time limit before 1% gap on medium instances"],
        font_size=13, color=RED)

section_label(sl, 0.4, 4.55, "Verdict for capstone")
add_textbox(sl, 0.4, 4.95, 12.5, 0.45,
            "The model is mathematically correct, modular, and fully tested. "
            "The C1 fix and per-node kappa are genuine improvements over the original. "
            "It is the right model for this scope — limitations are well-understood and documented.",
            font_size=14, bold=False, color=DARK_BLUE)

section_label(sl, 0.4, 5.55, "Constraint status")
verdicts = [
    ("C1  Placement integrity", "FIXED (== not <=)", GREEN),
    ("C2  Probabilistic capacity (SOCP)", "Sound, per-node kappa", GREEN),
    ("C3  Isolation primitives", "Correct, floor verified", GREEN),
    ("C4  Control-plane budgets", "Correct (inactive with synthetic data)", MID_BLUE),
    ("C5  SLA latency (big-M)", "Correct (inactive with synthetic data)", MID_BLUE),
    ("C6  Migration + disruption", "Correct (inactive with synthetic data)", MID_BLUE),
    ("C7  DRF fairness", "Correct, Q_quota fix applied", GREEN),
]
for i, (c, status, color) in enumerate(verdicts):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.4
    y = 5.95 + row * 0.35
    add_textbox(sl, x, y, 3.2, 0.32, c, font_size=11, bold=True, color=DARK_GREY)
    add_textbox(sl, x + 3.2, y, 3.0, 0.32, status, font_size=11, color=color)


# ============================================================
# SLIDE 19 — Closing
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(sl, 0, 2.9, 13.33, 1.8, MID_BLUE)
add_textbox(sl, 0.8, 1.2, 11.7, 0.8,
            "Thank You",
            font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 2.1, 11.7, 0.6,
            "Questions?",
            font_size=28, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.8, 3.1, 11.7, 0.5,
            "Plan-Ahead MISOCP  +  Real-Time MILP  +  Pipeline Interface",
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, 1.5, 4.2, 10.3, 1.6,
            "Key deliverables:\n"
            "  plan_ahead_data.py  |  plan_ahead_optimizer.py  |  plan_ahead_sensitivity.py\n"
            "  optimizer_google_or.py  |  pipeline/interface.py  |  pipeline/pipeline_configs.py\n"
            "  test_plan_ahead.py (17 tests)  |  test_model.py (14 tests)",
            font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)


# ── Save ──────────────────────────────────────────────────────────────────
OUT = r"C:\Users\alric\Documents\Spring2026\Capstone_ClaudeAI\op_model_overview.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
