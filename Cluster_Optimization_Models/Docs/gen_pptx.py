# -*- coding: utf-8 -*-
"""Generate optimization_model_overview.pptx"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).parent / "optimization_model_overview.pptx"

# ── Color palette ────────────────────────────────────────────────────────────
BG     = RGBColor(0x02, 0x06, 0x17)   # slate-950
CARD   = RGBColor(0x0f, 0x17, 0x2a)   # slate-900
BORDER = RGBColor(0x1e, 0x29, 0x3b)   # slate-800
WHITE  = RGBColor(0xff, 0xff, 0xff)
SLATE3 = RGBColor(0xcb, 0xd5, 0xe1)   # slate-300
SLATE4 = RGBColor(0x94, 0xa3, 0xb8)   # slate-400
SLATE5 = RGBColor(0x64, 0x74, 0x8b)   # slate-500
CYAN   = RGBColor(0x22, 0xd3, 0xee)   # cyan-400
AMBER  = RGBColor(0xfb, 0xbf, 0x24)   # amber-400
EMERALD= RGBColor(0x34, 0xd3, 0x99)   # emerald-400
PURPLE = RGBColor(0xa7, 0x8b, 0xfa)   # violet-400
SKY    = RGBColor(0x38, 0xbd, 0xf8)   # sky-400
ROSE   = RGBColor(0xfb, 0x71, 0x85)   # rose-400

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def bg(slide):
    bg_shape = slide.shapes.add_shape(1, 0, 0, W, H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG
    bg_shape.line.fill.background()


def textbox(slide, text, x, y, w, h,
            size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
            font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def card(slide, x, y, w, h, color=CARD):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.color.rgb = BORDER
    sh.line.width = Pt(0.75)
    return sh


def divider(slide, x, y, w, color=BORDER):
    ln = slide.shapes.add_shape(1, x, y, w, Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

# Accent bar
bar = sl.shapes.add_shape(1, 0, Inches(3.1), W, Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()

textbox(sl, "A Two-Part Optimization System",
        Inches(1), Inches(1.4), Inches(11.3), Inches(1.4),
        size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

textbox(sl, "Plan-Ahead Optimization  +  Real-Time Optimization",
        Inches(1), Inches(2.7), Inches(11.3), Inches(0.6),
        size=22, bold=False, color=CYAN, align=PP_ALIGN.CENTER)

textbox(sl, "The Optimization Layer of the Multi-Tenant Cluster Scheduling Pipeline",
        Inches(1), Inches(3.5), Inches(11.3), Inches(0.5),
        size=16, color=SLATE4, align=PP_ALIGN.CENTER)

textbox(sl, "Capstone Project  |  Spring 2026",
        Inches(1), Inches(6.6), Inches(11.3), Inches(0.4),
        size=13, color=SLATE5, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Two-Layer Pipeline
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

textbox(sl, "The Multi-Tenant Cluster Optimization Pipeline",
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
        size=28, bold=True, color=WHITE)
divider(sl, Inches(0.5), Inches(0.82), Inches(12.3), CYAN)

# Layer 1 card
card(sl, Inches(0.5), Inches(1.1), Inches(5.5), Inches(5.4))
textbox(sl, "Layer 1 — Prediction Layer",
        Inches(0.7), Inches(1.2), Inches(5.1), Inches(0.5),
        size=17, bold=True, color=AMBER)
lines1 = [
    "Processes historical cluster-usage traces",
    "  (Google Borg / Kubernetes telemetry)",
    "",
    "Outputs per-job predictions:",
    "  Pred_mem   max predicted memory usage",
    "  Pred_CPU   P95 predicted CPU peak",
    "",
    "Feeds both optimization models:",
    "  Plan-Ahead  uses aggregate demand u[i,h]",
    "  Real-Time   uses per-job Pred_mem, Pred_CPU",
]
textbox(sl, "\n".join(lines1),
        Inches(0.7), Inches(1.75), Inches(5.0), Inches(3.8),
        size=13, color=SLATE3)

# Arrow
textbox(sl, "►",
        Inches(6.15), Inches(3.4), Inches(0.6), Inches(0.5),
        size=28, color=CYAN, align=PP_ALIGN.CENTER)

# Layer 2 card
card(sl, Inches(6.8), Inches(1.1), Inches(6.0), Inches(5.4))
textbox(sl, "Layer 2 — Optimization Layer",
        Inches(7.0), Inches(1.2), Inches(5.6), Inches(0.5),
        size=17, bold=True, color=CYAN)

# Sub-card Plan-Ahead
card(sl, Inches(7.0), Inches(1.75), Inches(5.6), Inches(2.2), color=RGBColor(0x0a, 0x1a, 0x30))
textbox(sl, "Plan-Ahead (MILP / MISOCP)",
        Inches(7.15), Inches(1.82), Inches(5.2), Inches(0.4),
        size=13, bold=True, color=PURPLE)
textbox(sl, "Runs every H intervals (the horizon)\nAssigns machine pools and tenant groups\nfor the upcoming planning periods",
        Inches(7.15), Inches(2.25), Inches(5.2), Inches(1.4),
        size=12, color=SLATE3)

# Sub-card Real-Time
card(sl, Inches(7.0), Inches(4.05), Inches(5.6), Inches(2.1), color=RGBColor(0x0a, 0x1a, 0x30))
textbox(sl, "Real-Time (MILP)",
        Inches(7.15), Inches(4.12), Inches(5.2), Inches(0.4),
        size=13, bold=True, color=EMERALD)
textbox(sl, "Called once per tenant group per interval\nPlaces queued jobs onto assigned machines\nUpdates fairness weights (omega_t) each round",
        Inches(7.15), Inches(4.55), Inches(5.2), Inches(1.35),
        size=12, color=SLATE3)

textbox(sl, "Gurobi (WLS license)", Inches(7.0), Inches(3.4), Inches(2.6), Inches(0.35),
        size=11, color=SLATE5)
textbox(sl, "OR-Tools (open-source)", Inches(9.7), Inches(3.4), Inches(2.9), Inches(0.35),
        size=11, color=SLATE5, align=PP_ALIGN.RIGHT)

textbox(sl, "Cluster Manager routes pre-filtered jobs + machines to Real-Time each interval",
        Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.35),
        size=12, color=SLATE5, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Plan-Ahead Model Overview
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

textbox(sl, "Plan-Ahead Optimization  (MILP / MISOCP)",
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
        size=28, bold=True, color=WHITE)
divider(sl, Inches(0.5), Inches(0.82), Inches(12.3), PURPLE)

# Left: sets & decisions
card(sl, Inches(0.5), Inches(1.05), Inches(5.8), Inches(5.9))

textbox(sl, "Sets & Decisions",
        Inches(0.7), Inches(1.15), Inches(5.4), Inches(0.4),
        size=15, bold=True, color=PURPLE)

sets_txt = (
    "T = T_e (exclusive) u T_s (shared)\n"
    "M = M_a (always-on) u M_b (additional)\n"
    "H = planning periods (slots in horizon)\n\n"
    "Decision variables:\n"
    "  e[i,n,h]  1 if exclusive tenant i on machine n in period h\n"
    "            PER-PERIOD: adapts each period via feedback\n"
    "  z_on[n]   1 if additional machine n is activated\n"
    "  y[i,n,h]  1 if shared tenant i on machine n in period h\n"
    "  f[i,n,h]  capacity allocated to tenant i on machine n\n"
    "  sigma     min demand-satisfaction ratio (fairness)\n"
    "  t[n,h]    Cantelli safety buffer (MISOCP mode only)\n"
    "  mix[n,h]  1 if machine n has heavy + light tenant in h"
)
textbox(sl, sets_txt,
        Inches(0.7), Inches(1.6), Inches(5.4), Inches(4.2),
        size=11.5, color=SLATE3)

# Right: constraints & objective
card(sl, Inches(6.5), Inches(1.05), Inches(6.33), Inches(5.9))

textbox(sl, "Constraints & Objective",
        Inches(6.7), Inches(1.15), Inches(5.9), Inches(0.4),
        size=15, bold=True, color=PURPLE)

constr_txt = (
    "C_aa    z[n,h] = 1  for all n in M_a  (always-on)\n"
    "C_excl  each exclusive tenant -> exactly one machine\n"
    "C_sep   exclusive and shared machines do not overlap\n"
    "C_share each shared tenant assigned >= 1 machine / period\n\n"
    "Capacity (MILP):   Sum_i f[i,n,h]  <=  C_eff[n] * z[n,h]\n"
    "Capacity (MISOCP): above + kappa * t[n,h]  <=  C_eff[n] * z[n,h]\n"
    "  where  t[n,h]^2 >= Sum_i sigma2[i,h] * y[i,n,h]  (cone)\n"
    "  Cantelli: P(actual <= capacity) >= 1 - epsilon\n\n"
    "C_demand  Sum_n f[i,n,h] >= u[i,h]  (all shared, all periods)\n"
    "C_fair    sigma <= total_alloc / total_demand  (per tenant)\n\n"
    "Objective (minimize):\n"
    "  lam[0] * infra_cost\n"
    "  - lam[1] * sigma          (maximize fairness)\n"
    "  - lam[2] * mix_total      (reward heavy+light co-location)"
)
textbox(sl, constr_txt,
        Inches(6.7), Inches(1.6), Inches(5.9), Inches(5.0),
        size=11, color=SLATE3)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Real-Time Model Overview
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

textbox(sl, "Real-Time Optimization  (MILP per group per interval)",
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
        size=28, bold=True, color=WHITE)
divider(sl, Inches(0.5), Inches(0.82), Inches(12.3), EMERALD)

# Left
card(sl, Inches(0.5), Inches(1.05), Inches(5.8), Inches(5.9))
textbox(sl, "Model & Inputs",
        Inches(0.7), Inches(1.15), Inches(5.4), Inches(0.4),
        size=15, bold=True, color=EMERALD)

rt_left = (
    "Solver: OR-Tools  (CBC / GLOP / SCIP)\n"
    "Called: once per tenant group per interval\n"
    "Inputs: pre-filtered jobs + machines from Cluster Manager\n\n"
    "Decision variable:\n"
    "  x[j,n]  1 if job j placed on machine n\n\n"
    "Derived node quantities:\n"
    "  v_bar_n   SLA violation rate (last K intervals)\n"
    "  M_cap_n   schedulable capacity = RAM - OS_tax - theta\n"
    "  M_avail_n = M_cap_n - used_n\n"
    "  M_eff_n   = max(0, M_avail_n * (1 - v_bar_n))\n"
    "  u_mem_n   = 1 + used_n / M_cap_n  (utilization weight)\n\n"
    "Delay weight (fairness feedback):\n"
    "  omega_t = 1 + max(0, (W_bar_t - W_bar) / max(1, W_bar))\n"
    "  Tenants waiting longer than average get omega_t > 1"
)
textbox(sl, rt_left,
        Inches(0.7), Inches(1.6), Inches(5.4), Inches(4.2),
        size=11.5, color=SLATE3)

# Right
card(sl, Inches(6.5), Inches(1.05), Inches(6.33), Inches(5.9))
textbox(sl, "Objective & Constraints",
        Inches(6.7), Inches(1.15), Inches(5.9), Inches(0.4),
        size=15, bold=True, color=EMERALD)

rt_right = (
    "Objective (maximize weighted memory placement):\n\n"
    "  max Z = Sum_{j,n}  omega_t(j)  *  Pred_mem_j\n"
    "                  *  u_mem_n  *  sigma_consolid_n  *  x[j,n]\n\n"
    "  omega_t(j)       delay weight (fairness boost)\n"
    "  Pred_mem_j       predicted memory of job j\n"
    "  u_mem_n          utilization weight (consolidation)\n"
    "  sigma_consolid_n node consolidation bias\n\n"
    "Constraints:\n"
    "  C1: Sum_n x[j,n] <= 1      (one machine per job)\n"
    "  C2: Sum_j Pred_mem_j * x[j,n] <= M_eff_n\n"
    "                             (effective memory budget)\n"
    "  C3: x[j,n] in {0,1}       (binary placement)\n"
    "  C4: x[j,n] = 0  if Pred_CPU_j > CPU_n\n"
    "                             (CPU fitment, pre-filtered)\n\n"
    "Unplaced jobs stay in queue; tenant W_bar_t bumped by\n"
    "one interval -- raising omega_t next round."
)
textbox(sl, rt_right,
        Inches(6.7), Inches(1.6), Inches(5.9), Inches(5.0),
        size=11.5, color=SLATE3)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Simulation Application
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

textbox(sl, "Simulation Application",
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
        size=28, bold=True, color=WHITE)
divider(sl, Inches(0.5), Inches(0.82), Inches(12.3), SKY)

textbox(sl, "Interactive visualization of the Multi-Tenant Cluster Optimization Pipeline running in real time",
        Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.45),
        size=15, color=SLATE4, align=PP_ALIGN.CENTER)

# Three feature cards
cw = Inches(3.8)
ch = Inches(4.4)
cy = Inches(1.6)

card(sl, Inches(0.5), cy, cw, ch)
textbox(sl, "Live Cluster View",
        Inches(0.7), cy + Inches(0.1), cw - Inches(0.3), Inches(0.4),
        size=15, bold=True, color=SKY)
textbox(sl, (
    "Job queue with tenant color coding\n"
    "Node cards showing RAM usage,\n"
    "  SLA violation rate, and job count\n"
    "Memory utilization wave chart\n"
    "Auto-pause when Plan-Ahead fires\n\n"
    "Metrics displayed:\n"
    "  Cap Util   used / physical RAM\n"
    "  Eff Util   used / M_cap (schedulable)\n"
    "  Act Util   eff over active nodes only"
), Inches(0.7), cy + Inches(0.55), cw - Inches(0.3), Inches(3.5),
    size=12, color=SLATE3)

card(sl, Inches(4.77), cy, cw, ch)
textbox(sl, "Plan-Ahead Gantt",
        Inches(4.97), cy + Inches(0.1), cw - Inches(0.3), Inches(0.4),
        size=15, bold=True, color=PURPLE)
textbox(sl, (
    "Runs every H intervals (configurable)\n"
    "Gantt chart showing machine groups\n"
    "  per tenant per planning period\n"
    "Highlights active period with marker\n"
    "Exclusive vs shared tenant coloring\n\n"
    "Gurobi (MILP or MISOCP) when licensed\n"
    "Falls back to mock planner otherwise\n"
    "On-demand via Plan Ahead button"
), Inches(4.97), cy + Inches(0.55), cw - Inches(0.3), Inches(3.5),
    size=12, color=SLATE3)

card(sl, Inches(9.03), cy, cw, ch)
textbox(sl, "Configurable Parameters",
        Inches(9.23), cy + Inches(0.1), cw - Inches(0.3), Inches(0.4),
        size=15, bold=True, color=AMBER)
textbox(sl, (
    "Topology: nodes (default 20), tenants,\n"
    "  always-on (4), RAM [64-128 GB], CPU [4-16]\n"
    "Workload: jobs/interval [0-10],\n"
    "  job RAM/CPU, spike prob (5%),\n"
    "  job lifetime [4-60 s]\n"
    "Scheduler: K window, safety buffer\n"
    "Plan-Ahead: horizon (24 steps),\n"
    "  period width (6 steps), exclusivity,\n"
    "  MILP vs MISOCP, feedback refs\n"
    "  (W_ref=1s, q_ref=10, gamma=0.3)\n\n"
    "All changes staged, applied on Reset"
), Inches(9.23), cy + Inches(0.55), cw - Inches(0.3), Inches(3.5),
    size=12, color=SLATE3)

textbox(sl, "Backend: FastAPI + Python   |   Frontend: React + TypeScript + Tailwind   |   Solvers: Gurobi (plan-ahead) + OR-Tools (real-time)",
        Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
        size=11, color=SLATE5, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Feedback & Sensitivity
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

textbox(sl, "Feedback Integration & Sensitivity Analysis",
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
        size=28, bold=True, color=WHITE)
divider(sl, Inches(0.5), Inches(0.82), Inches(12.3), ROSE)

# Feedback card
card(sl, Inches(0.5), Inches(1.05), Inches(5.8), Inches(5.5))
textbox(sl, "Feedback Integration",
        Inches(0.7), Inches(1.15), Inches(5.4), Inches(0.4),
        size=15, bold=True, color=ROSE)

fb_txt = (
    "Plan-Ahead receives 3 signals from Real-Time:\n\n"
    "1. SLA capacity reduction (per machine):\n"
    "   C_eff[n] = C[n] * (1 - alpha * v_bar_n)\n"
    "   Machines with frequent violations get smaller\n"
    "   effective capacity in the next plan solve\n\n"
    "2. Wait-time demand inflation (ALL tenants):\n"
    "   wait_scale = 1 + beta * min(2, W_bar_i / W_ref)\n\n"
    "3. Queue-size demand inflation (ALL tenants):\n"
    "   queue_scale = 1 + gamma * min(2, q_i / q_ref)\n"
    "   Combined: u[i,h] = u_raw[i,h] * min(3, wait * queue)\n"
    "   Exclusive tenants also inflate -> more machines assigned\n\n"
    "Cantelli safety buffer:\n"
    "   kappa = sqrt((1-eps)/eps); eps=0.10 -> kappa=3.0\n"
    "   t[n,h]^2 >= Sum_i sigma2[i,h] * y[i,n,h]\n"
    "   P(actual <= C_eff[n]) >= 1 - eps = 90% safety"
)
textbox(sl, fb_txt,
        Inches(0.7), Inches(1.6), Inches(5.4), Inches(4.6),
        size=11.5, color=SLATE3)

# Sensitivity card — actual pipeline results
card(sl, Inches(6.5), Inches(1.05), Inches(6.33), Inches(5.5))
textbox(sl, "Pipeline Sensitivity Results",
        Inches(6.7), Inches(1.15), Inches(5.9), Inches(0.4),
        size=15, bold=True, color=ROSE)

sens_txt = (
    "Baseline: 8-node cluster, 8 GB/node, 4 tenants\n"
    "All sweeps: 30 intervals  |  OR-Tools + Gurobi mock\n\n"
    "Arrival rate (J = 1-8 jobs/interval):\n"
    "  100% placed across all rates; queue = 0\n"
    "  RT solver: 1 ms (J=1) -> 15 ms (J=8)\n"
    "  -> Sub-real-time even at peak test load\n\n"
    "Tenant count (T = 2-10):\n"
    "  PA vars: 200 (T=2) -> 840 (T=10); solve 10 ms\n"
    "  -> Linear scaling; T <= 15 feasible per group\n\n"
    "Exclusivity (0-3 exclusive out of 6):\n"
    "  PA vars decrease as excl rises (fewer shared combos)\n"
    "  -> Keep exclusive < 25% of tenants\n\n"
    "Node count (N = 3-20):\n"
    "  PA vars: 135 (N=3) -> 900 (N=20); all 100% placed\n"
    "  -> Rule: 1 node per ~2-3 concurrent sustained jobs\n\n"
    "Job lifetime (5-120 s):\n"
    "  eff_mem% stays 0% up to 60 s; rises to 5.5% at 120 s\n"
    "  -> Sweet spot: 15-30 s; set horizon >= 2x max lifetime"
)
textbox(sl, sens_txt,
        Inches(6.7), Inches(1.6), Inches(5.9), Inches(4.6),
        size=11, color=SLATE3)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Scalability & Model Summary
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

textbox(sl, "Scalability Analysis & Model Summary",
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
        size=28, bold=True, color=WHITE)
divider(sl, Inches(0.5), Inches(0.82), Inches(12.3), AMBER)

# Left card — model what-does-it-do summary
card(sl, Inches(0.5), Inches(1.05), Inches(5.8), Inches(5.7))
textbox(sl, "What Each Model Does",
        Inches(0.7), Inches(1.15), Inches(5.4), Inches(0.4),
        size=15, bold=True, color=AMBER)

model_txt = (
    "PLAN-AHEAD  (MILP / MISOCP, Gurobi)\n"
    "  Runs once per planning horizon (every H intervals).\n"
    "  Decides WHICH machines are active and WHICH tenant\n"
    "  groups get access to them for the next H intervals.\n"
    "  Exclusive tenants get private machine pools per period;\n"
    "  shared tenants compete within a common pool.\n"
    "  Demand is inflated by feedback (wait time + queue size)\n"
    "  so the plan scales up capacity before queues spike.\n"
    "  Cantelli cone ensures P(OOM) < 10% (MISOCP mode).\n\n"
    "REAL-TIME  (MILP, OR-Tools)\n"
    "  Runs once per tenant group per interval (every second).\n"
    "  Decides WHICH queued jobs go on WHICH machines NOW.\n"
    "  Uses pre-filtered machines from the Plan-Ahead output.\n"
    "  Fairness weight omega_t boosts tenants that have been\n"
    "  waiting longer than the cluster average, preventing\n"
    "  starvation even as new high-demand jobs arrive.\n"
    "  Unplaced jobs stay in queue; omega_t rises next round.\n\n"
    "TOGETHER: Plan-Ahead manages the allocation horizon;\n"
    "Real-Time maximizes throughput within that allocation."
)
textbox(sl, model_txt,
        Inches(0.7), Inches(1.6), Inches(5.4), Inches(5.0),
        size=11, color=SLATE3)

# Right card — scalability rules
card(sl, Inches(6.5), Inches(1.05), Inches(6.33), Inches(5.7))
textbox(sl, "Scalability Rules of Thumb",
        Inches(6.7), Inches(1.15), Inches(5.9), Inches(0.4),
        size=15, bold=True, color=AMBER)

scale_txt = (
    "Real-Time MILP  [O(J x N) binary vars]\n"
    "  < 10 ms  at J=1-8, N=8    (tested)\n"
    "  < 100 ms at J=50, N=20\n"
    "  ~ 500 ms at J=200, N=50  (time-limit applied)\n"
    "  VIABLE: live per-interval scheduler for J <= 50, N <= 20\n\n"
    "Plan-Ahead MILP/MISOCP  [O(T x N x H) vars]\n"
    "  T=4, N=8, H=5 -> 360 vars, < 5 s solve (tested)\n"
    "  T=10, N=20, H=6 -> ~840 vars, ~ 15 s solve\n"
    "  T=25, N=50, H=8 -> > 60 s (use hierarchical partition)\n"
    "  VIABLE: periodic planner (every 24 intervals).\n"
    "  Larger deployments need tenant group decomposition.\n\n"
    "Capacity sizing:\n"
    "  Add 1 node per ~2-3 sustained concurrent jobs\n"
    "  (at avg 1.3 GB/job x 8 GB nodes)\n\n"
    "Exclusivity budget:\n"
    "  Keep exclusive count < 25% of total tenants\n"
    "  Above that, shared pool overflow rises sharply\n\n"
    "Lifetime vs horizon:\n"
    "  Plan-Ahead horizon >= 2x max job lifetime\n"
    "  Sweet spot: 15-30 s lifetime, horizon = 24 steps"
)
textbox(sl, scale_txt,
        Inches(6.7), Inches(1.6), Inches(5.9), Inches(5.0),
        size=11, color=SLATE3)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Thank You
# ════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

bar2 = sl.shapes.add_shape(1, 0, Inches(3.6), W, Inches(0.04))
bar2.fill.solid(); bar2.fill.fore_color.rgb = CYAN; bar2.line.fill.background()

textbox(sl, "Thank You",
        Inches(1), Inches(1.8), Inches(11.3), Inches(1.0),
        size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

textbox(sl, "A Two-Part Optimization System for Multi-Tenant Cluster Scheduling",
        Inches(1), Inches(3.8), Inches(11.3), Inches(0.55),
        size=20, color=CYAN, align=PP_ALIGN.CENTER)

textbox(sl, "Gurobi MILP/MISOCP  (Plan-Ahead)   +   OR-Tools MILP  (Real-Time)   +   FastAPI / React Simulation",
        Inches(1), Inches(4.6), Inches(11.3), Inches(0.4),
        size=14, color=SLATE4, align=PP_ALIGN.CENTER)


prs.save(str(OUT))
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
