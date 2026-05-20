"""
generate_slides.py
──────────────────
Generates plan_ahead_refactoring.pptx using python-pptx.
Run from any directory:
    python Cluster_Optimization_Models/Docs/generate_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)
MID_BLUE    = RGBColor(0x2A, 0x5A, 0x8C)
LIGHT_BLUE  = RGBColor(0xE8, 0xF0, 0xFE)
ACCENT      = RGBColor(0x0D, 0x6E, 0xFD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MED_GRAY    = RGBColor(0xCC, 0xD6, 0xE0)
GREEN_BG    = RGBColor(0xE6, 0xF4, 0xEA)
GREEN_FG    = RGBColor(0x1E, 0x7E, 0x34)
RED_BG      = RGBColor(0xFC, 0xE8, 0xE6)
RED_FG      = RGBColor(0xC6, 0x28, 0x28)

# Slide dimensions: widescreen 16:9
W = Inches(13.33)
H = Inches(7.5)

FONT = "Calibri"


# ── Low-level helpers ──────────────────────────────────────────────────────

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             wrap=True, font_name=FONT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=16, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def header_bar(slide, title, subtitle=None):
    """Dark blue header bar across top."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=DARK_BLUE)
    add_text(slide, title, 0.35, 0.08, 11.5, 0.7,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.68, 10, 0.38,
                 font_size=16, italic=True, color=MED_GRAY, align=PP_ALIGN.LEFT)
    # thin accent line under header
    add_rect(slide, 0, 1.1, 13.33, 0.04, fill_rgb=ACCENT)


def footer(slide, prs, num):
    """Page number + branding footer."""
    add_rect(slide, 0, 7.18, 13.33, 0.32, fill_rgb=DARK_BLUE)
    add_text(slide, "Alrick Grandison  |  Capstone Project, Spring 2026",
             0.3, 7.2, 9, 0.28, font_size=12, color=MED_GRAY)
    add_text(slide, str(num), 12.6, 7.2, 0.5, 0.28,
             font_size=12, color=MED_GRAY, align=PP_ALIGN.RIGHT)


def bullet_box(slide, items, x, y, w, h,
               font_size=16, title=None, title_size=17,
               bg=None, indent="  •  "):
    """A bordered box with optional title and bullet lines."""
    if bg:
        add_rect(slide, x, y, w, h, fill_rgb=bg,
                 line_rgb=MED_GRAY, line_width_pt=0.5)
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.12), Inches(y + 0.12),
        Inches(w - 0.24), Inches(h - 0.24)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name = FONT
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = DARK_BLUE
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{indent}{item}"
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = NEAR_BLACK
    return txBox


def two_col_table(slide, headers, rows, x, y, w, h,
                  col_widths=None, font_size=15, header_font_size=15):
    """Render a simple table as stacked rectangles."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols
    n_rows = len(rows)
    row_h = h / (n_rows + 1)

    # header row
    cx = x
    for ci, (hdr, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, cx, y, cw, row_h, fill_rgb=DARK_BLUE)
        add_text(slide, hdr, cx + 0.05, y + 0.03, cw - 0.1, row_h - 0.06,
                 font_size=header_font_size, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True)
        cx += cw

    # data rows
    for ri, row in enumerate(rows):
        ry = y + row_h * (ri + 1)
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        cx = x
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, cx, ry, cw, row_h, fill_rgb=bg,
                     line_rgb=MED_GRAY, line_width_pt=0.3)
            add_text(slide, str(cell), cx + 0.05, ry + 0.03, cw - 0.1, row_h - 0.06,
                     font_size=font_size, color=NEAR_BLACK,
                     align=PP_ALIGN.LEFT, wrap=True)
            cx += cw


def callout(slide, text, x, y, w, h, bg=LIGHT_BLUE, fg=MID_BLUE, font_size=15):
    add_rect(slide, x, y, w, h, fill_rgb=bg, line_rgb=MID_BLUE, line_width_pt=1)
    add_rect(slide, x, y, 0.07, h, fill_rgb=MID_BLUE)
    add_text(slide, text, x + 0.18, y + 0.08, w - 0.3, h - 0.16,
             font_size=font_size, color=fg, wrap=True, italic=True)


def status_badge(slide, text, x, y, ok=True):
    col = GREEN_FG if ok else RED_FG
    bg  = GREEN_BG if ok else RED_BG
    add_rect(slide, x, y, 0.9, 0.28, fill_rgb=bg, line_rgb=col, line_width_pt=0.5)
    add_text(slide, text, x + 0.06, y + 0.03, 0.8, 0.22,
             font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)


def section_divider(prs, part_num, part_title, subtitle):
    """Full-bleed dark section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.title  # suppress default
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
    add_rect(slide, 0, 3.1, 13.33, 0.06, fill_rgb=ACCENT)
    add_text(slide, f"PART {part_num}", 0, 2.2, 13.33, 0.5,
             font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, part_title, 0, 2.7, 13.33, 0.9,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 0, 3.6, 13.33, 0.5,
             font_size=20, italic=True, color=MED_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Alrick Grandison  |  Capstone Project, Spring 2026",
             0, 6.9, 13.33, 0.4, font_size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)
    return slide


# ── Slide builders ─────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
    add_rect(slide, 0, 3.5, 13.33, 0.06, fill_rgb=ACCENT)
    add_text(slide, "Plan-Ahead Model Refactoring",
             0, 1.5, 13.33, 0.9, font_size=42, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Justification, Design Decisions & Optimization Layer Status",
             0, 2.5, 13.33, 0.55, font_size=22, italic=True,
             color=MED_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Alrick Grandison  |  Capstone Project, Spring 2026",
             0, 4.0, 13.33, 0.45, font_size=18,
             color=MED_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Multi-Tenant Cluster Scheduler — Optimization Layer",
             0, 4.6, 13.33, 0.4, font_size=15,
             color=_rgb(0x88, 0xA8, 0xC8), align=PP_ALIGN.CENTER)


def slide_agenda(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Agenda")
    footer(slide, prs, num)

    # Left column
    bullet_box(slide,
        ["What the original model was",
         "Problems discovered during simulation",
         "Change 1: Workloads → Tenant Usage Profiles",
         "Change 2: Hard Access → Priority Hints",
         "Change 3: SOCP uncertainty input redesigned",
         "What was removed (migration, isolation, McCormick)",
         "Complexity comparison — model remains MISOCP"],
        x=0.3, y=1.3, w=6.0, h=4.5,
        bg=LIGHT_GRAY, title="Part 1 — Plan-Ahead Refactoring",
        font_size=15)

    # Right column
    bullet_box(slide,
        ["System architecture overview",
         "Plan-Ahead: current state",
         "Real-Time: current state",
         "Pipeline: how the layers connect",
         "Simulation: current state",
         "What’s next",
         "Future work"],
        x=6.7, y=1.3, w=6.3, h=4.5,
        bg=LIGHT_BLUE, title="Part 2 — Optimization Layer Status",
        font_size=15)


def slide_original_model(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "The Original Model", "What we started with")
    footer(slide, prs, num)

    two_col_table(slide,
        ["Component", "Description"],
        [
            ["Scheduling unit",   "Individual workloads  x_{i,j,n,t} in {0,1}  — per-job binary placement variable"],
            ["Capacity (SOCP)",   "Cantelli cone on per-workload covariance Σ_r (Cholesky decomposition L_r)"],
            ["Access control",    "Hard block: x_{j,n} = 0  if node n not in A_{t(j)}"],
            ["Isolation",         "gVisor / Kata selection  w_{i,j,k,t}  with McCormick linearization (ξ variables)"],
            ["Migration",         "Explicit migration indicators  m_{i,j,n,t}  across planning periods"],
            ["Fairness",          "DRF max-min on declared per-job demand  d_{i,j,r}"],
        ],
        x=0.3, y=1.25, w=12.7, h=4.3,
        col_widths=[2.4, 10.3], font_size=14, header_font_size=15)

    callout(slide,
        "The model was mathematically complete on paper, "
        "but its inputs required data that cannot be known at planning time.",
        x=0.3, y=5.7, w=12.7, h=0.9, font_size=15)


def slide_problems(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Problems Found During Simulation")
    footer(slide, prs, num)

    # Box 1
    add_rect(slide, 0.3, 1.25, 12.73, 1.5, fill_rgb=LIGHT_GRAY,
             line_rgb=MED_GRAY, line_width_pt=0.5)
    add_text(slide, "1  —  Infeasible Input Requirements", 0.5, 1.3, 10, 0.38,
             font_size=16, bold=True, color=DARK_BLUE)
    add_text(slide,
        "The SOCP cone required a Cholesky decomposition of Σ_r — "
        "the joint covariance matrix of all individual workloads across the horizon. "
        "This is unknowable in advance; workloads are only known when jobs arrive.",
        0.5, 1.65, 12.3, 0.95, font_size=14, wrap=True)

    # Box 2
    add_rect(slide, 0.3, 2.9, 12.73, 1.65, fill_rgb=LIGHT_GRAY,
             line_rgb=MED_GRAY, line_width_pt=0.5)
    add_text(slide, "2  —  Feasibility Failures & Queue Starvation", 0.5, 2.95, 10, 0.38,
             font_size=16, bold=True, color=DARK_BLUE)
    add_text(slide,
        "With x_{j,n} = 0 outside the plan-ahead set, jobs were blocked from idle nodes "
        "while their “allowed” nodes were saturated.\n"
        "Little’s Law: blocking a job from 4 of 5 nodes multiplies expected wait time proportionally.\n"
        "Simulation confirmed: queue backlogs persisted even with idle capacity available.",
        0.5, 3.3, 12.3, 1.1, font_size=14, wrap=True)

    # Box 3
    add_rect(slide, 0.3, 4.7, 12.73, 1.35, fill_rgb=LIGHT_GRAY,
             line_rgb=MED_GRAY, line_width_pt=0.5)
    add_text(slide, "3  —  Over-Engineered Variables", 0.5, 4.75, 10, 0.38,
             font_size=16, bold=True, color=DARK_BLUE)
    add_text(slide,
        "Migration and isolation-primitive variables each scale as O(T · W · N · K · H). "
        "Container isolation choice and migration decisions are real-time concerns — "
        "they cannot be meaningfully set at planning granularity.",
        0.5, 5.1, 12.3, 0.85, font_size=14, wrap=True)

    callout(slide, "Key insight: the original design modelled the wrong level of detail for a plan-ahead horizon.",
            x=0.3, y=6.2, w=12.73, h=0.65, font_size=14)


def slide_change1(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Change 1: Workloads → Tenant Usage Profiles",
               "What the planning unit is and why it changed")
    footer(slide, prs, num)

    two_col_table(slide,
        ["", "Original", "Refactored"],
        [
            ["Planning unit",        "x_{i,j,n,t} in {0,1}  per workload",
                                     "u_{i,h} >= 0  — aggregate usage profile per tenant per period"],
            ["What is predicted",    "Which specific jobs tenant i will submit",
                                     "How much resource tenant i uses in total during period h"],
            ["Predictability",       "Very low — depends on user behaviour & pipelines",
                                     "High — aggregate demand forecastable at 15-min to 4-hr granularity"],
            ["Variable count",       "O(T · W · N · H)  (per-workload)",
                                     "O(T · N · H)  (per-tenant)"],
        ],
        x=0.3, y=1.25, w=12.7, h=3.5,
        col_widths=[2.5, 4.8, 5.4], font_size=13, header_font_size=14)

    add_text(slide, "Why aggregate demand is forecastable", 0.3, 4.9, 12.7, 0.38,
             font_size=16, bold=True, color=DARK_BLUE)
    add_text(slide,
        "Google cluster traces v3 show aggregate per-tenant resource consumption is predictable "
        "at 15-min to 4-hour granularity using standard time-series methods (ARIMA, LSTM). "
        "Individual job arrivals are not.",
        0.3, 5.28, 12.7, 0.7, font_size=14, wrap=True)

    callout(slide,
        "u_{i,h} = total expected resource usage of tenant i in period h"
        "  —  this is the clean, well-defined input the prediction team produces.",
        x=0.3, y=6.1, w=12.7, h=0.72, font_size=14)


def slide_change2(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Change 2: Hard Access → Priority Hints",
               "Replacing blocking constraints with objective boosts")
    footer(slide, prs, num)

    # Left: the problem
    add_rect(slide, 0.3, 1.25, 5.8, 2.9, fill_rgb=RED_BG,
             line_rgb=RED_FG, line_width_pt=0.8)
    add_text(slide, "✗  Original: Hard Block", 0.5, 1.3, 5.4, 0.4,
             font_size=16, bold=True, color=RED_FG)
    add_text(slide, "x_{j,n} = 0   if  n ∉ A_{t(j)}",
             0.5, 1.72, 5.4, 0.45, font_size=15, bold=True, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_text(slide,
        "Feasibility fragility: if the plan is stale or a node fails, "
        "legitimate jobs are unnecessarily rejected.\n\n"
        "Queue starvation: blocked tenants wait indefinitely while other "
        "nodes sit idle — violates fairness.",
        0.5, 2.22, 5.4, 1.8, font_size=13, wrap=True)

    # Right: the fix
    add_rect(slide, 7.2, 1.25, 5.83, 2.9, fill_rgb=GREEN_BG,
             line_rgb=GREEN_FG, line_width_pt=0.8)
    add_text(slide, "✓  Refactored: Priority Boost", 7.4, 1.3, 5.4, 0.4,
             font_size=16, bold=True, color=GREEN_FG)
    add_text(slide, "b_{t(j),n} = 2.0  if n in priority_set[t(j)]",
             7.4, 1.72, 5.4, 0.3, font_size=13, bold=True, color=NEAR_BLACK)
    add_text(slide, "b_{t(j),n} = 1.0  otherwise",
             7.4, 2.02, 5.4, 0.3, font_size=13, bold=True, color=NEAR_BLACK)
    add_text(slide,
        "Jobs are never unnecessarily blocked.\n"
        "Plan-ahead guidance is respected when capacity allows.\n"
        "System degrades gracefully when the plan is stale.",
        7.4, 2.45, 5.4, 1.5, font_size=13, wrap=True)

    # Arrow
    add_text(slide, "➡", 6.15, 2.4, 0.9, 0.6, font_size=30,
             color=DARK_BLUE, align=PP_ALIGN.CENTER)

    callout(slide,
        "No node is off-limits. The plan-ahead guides placements; it does not dictate them. "
        "Simulation confirmed: switching to priority-boost eliminated queue backlogs without "
        "sacrificing the planning signal.",
        x=0.3, y=4.35, w=12.7, h=0.9, font_size=14)

    # Little's Law note
    add_rect(slide, 0.3, 5.4, 12.7, 0.9, fill_rgb=LIGHT_GRAY,
             line_rgb=MED_GRAY, line_width_pt=0.5)
    add_text(slide,
        "Little’s Law: blocking a job from 4 of 5 nodes (default 10% exclusivity = 1 allowed node) "
        "multiplies expected wait time by 5× compared to free placement.",
        0.5, 5.47, 12.3, 0.75, font_size=13, wrap=True)


def slide_change3(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Change 3: SOCP Uncertainty Input Redesigned",
               "Same cone structure, different and obtainable uncertainty source")
    footer(slide, prs, num)

    # Old SOCP
    add_rect(slide, 0.3, 1.25, 5.9, 3.3, fill_rgb=LIGHT_GRAY,
             line_rgb=MED_GRAY, line_width_pt=0.5)
    add_text(slide, "Original C2 — Per-Workload Cholesky",
             0.5, 1.3, 5.5, 0.38, font_size=15, bold=True, color=DARK_BLUE)
    add_text(slide,
        "Mean load\n"
        "+ κ_n · ‖ L_r · ξ_{n,r,t} ‖₂\n"
        "≤  C_{n,r} · z_{n,t}",
        0.5, 1.72, 5.5, 1.0, font_size=14, wrap=True)
    add_text(slide,
        "Requires:  full workload covariance matrix Σ_r,\n"
        "Cholesky factor L_r  (Cholesky of Σ_r)\n\n"
        "➡ unknowable at planning time",
        0.5, 2.78, 5.5, 1.6, font_size=13, wrap=True, color=RED_FG)

    # Arrow
    add_text(slide, "➡", 6.35, 2.55, 0.6, 0.6, font_size=30,
             color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # New SOCP
    add_rect(slide, 7.1, 1.25, 5.93, 3.3, fill_rgb=GREEN_BG,
             line_rgb=GREEN_FG, line_width_pt=0.8)
    add_text(slide, "Refactored C1a + C1b — Per-Tenant Variance",
             7.3, 1.3, 5.6, 0.38, font_size=15, bold=True, color=GREEN_FG)
    add_text(slide, "C1a (linear):   sum_i f_{i,n,h} + κ · t_{n,h} ≤ C_n · z_{n,h}",
             7.3, 1.72, 5.6, 0.38, font_size=13, wrap=True)
    add_text(slide, "C1b (cone):    sum_i σ²_{i,h} · y_{i,n,h} ≤ t_{n,h}²",
             7.3, 2.12, 5.6, 0.38, font_size=13, wrap=True)
    add_text(slide, "where  σ²_{i,h} = (sigma_frac × u_{i,h})²",
             7.3, 2.55, 5.6, 0.38, font_size=13, italic=True, wrap=True)
    add_text(slide,
        "Requires:  sigma_frac (one operator scalar)\n"
        "+ u_{i,h} (prediction team output)\n\n"
        "➡ estimable from cluster usage traces",
        7.3, 2.98, 5.6, 1.4, font_size=13, wrap=True, color=GREEN_FG)

    callout(slide,
        "Same probabilistic guarantee:  P[ actual usage ≤ C_n · z_{n,h} ] ≥ 1 − ε\n"
        "Cantelli factor: κ = sqrt((1−ε)/ε)  —  ε=0.10 → κ=3.00  (90% coverage)",
        x=0.3, y=4.7, w=12.7, h=0.9, font_size=14)

    two_col_table(slide,
        ["", "Original", "Refactored"],
        [
            ["Uncertainty source", "Per-workload covariance Σ_r", "Per-tenant usage variance σ²_{i,h}"],
            ["Data required",      "Joint workload distribution (unknowable)", "sigma_frac × u_{i,h} (estimable)"],
            ["Cone variable",      "Tied to Cholesky of Σ_r", "Single slack t_{n,h} per node per period"],
        ],
        x=0.3, y=5.75, w=12.7, h=1.45,
        col_widths=[2.4, 5.0, 5.3], font_size=13, header_font_size=13)


def slide_misocp_preserved(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "The Model Remains a MISOCP — Why This Matters")
    footer(slide, prs, num)

    # Big green banner
    add_rect(slide, 0.3, 1.25, 12.73, 0.7, fill_rgb=GREEN_BG,
             line_rgb=GREEN_FG, line_width_pt=1.5)
    add_text(slide,
        "✓  The refactored model is still a Mixed-Integer Second-Order Cone Program (MISOCP). "
        "The cone constraint is preserved. What changed is the uncertainty source driving it.",
        0.5, 1.3, 12.3, 0.6, font_size=15, bold=True, color=GREEN_FG, wrap=True)

    two_col_table(slide,
        ["", "Original", "Refactored"],
        [
            ["Cone expression",
             "‖ L_r · ξ_{n,r,t} ‖₂ ≤ soc_aux",
             "sum_i σ²_{i,h} · y_{i,n,h} ≤ t_{n,h}²"],
            ["Uncertainty source",
             "Per-workload covariance Σ_r",
             "Per-tenant usage variance σ²_{i,h}"],
            ["Data required",
             "Joint workload distribution (unknowable in advance)",
             "sigma_frac × u_{i,h}  (estimable from usage traces)"],
            ["Probabilistic guarantee",
             "P[usage ≤ C] ≥ 1−ε",
             "P[usage ≤ C] ≥ 1−ε  (same)"],
            ["Model class",
             "MISOCP",
             "MISOCP"],
        ],
        x=0.3, y=2.1, w=12.7, h=3.2,
        col_widths=[2.5, 5.0, 5.2], font_size=13, header_font_size=14)

    add_rect(slide, 0.3, 5.45, 12.73, 1.15, fill_rgb=LIGHT_BLUE,
             line_rgb=MID_BLUE, line_width_pt=0.5)
    add_text(slide, "Simulation-speed MILP mode  (use_socp=False)",
             0.5, 5.5, 12.0, 0.38, font_size=15, bold=True, color=DARK_BLUE)
    add_text(slide,
        "A plain MILP fallback (drops C1b and t_{n,h}) is available for interactive simulation "
        "where the plan-ahead fires every 50 steps. "
        "This is NOT the primary formulation and does not change the model’s MISOCP identity.",
        0.5, 5.86, 12.2, 0.65, font_size=13, wrap=True, color=NEAR_BLACK)


def slide_removed(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "What Was Removed and Why")
    footer(slide, prs, num)

    items = [
        ("Isolation primitives  w_{i,j,k,t}  and McCormick auxiliaries  ξ_{i,j,n,k,t}",
         "Container isolation (gVisor / Kata) is a runtime decision made at placement time. "
         "It cannot be meaningfully set at planning granularity. "
         "Removing these eliminates McCormick linearizations and C3/C4 control-plane constraints."),
        ("Migration variables  m_{i,j,n,t}",
         "In a priority-hint model, the plan-ahead does not place individual jobs. "
         "Migration is a real-time concept — 'workload j moves from node n to n’ in period h' "
         "has no meaning when the plan only signals priority. Removes C6 entirely."),
        ("DRF fairness on declared demand  d_{i,j,r}",
         "Per-job declared demand d_{i,j,r} is not available at planning time. "
         "Fairness is now expressed over aggregate allocation vs aggregate demand (C4), "
         "which is computable from u_{i,h}."),
    ]

    y = 1.25
    for title, desc in items:
        add_rect(slide, 0.3, y, 12.73, 1.55, fill_rgb=LIGHT_GRAY,
                 line_rgb=MED_GRAY, line_width_pt=0.4)
        add_rect(slide, 0.3, y, 0.07, 1.55, fill_rgb=DARK_BLUE)
        add_text(slide, title, 0.5, y + 0.07, 12.2, 0.4,
                 font_size=14, bold=True, color=DARK_BLUE, wrap=True)
        add_text(slide, desc, 0.5, y + 0.5, 12.2, 0.95,
                 font_size=13, wrap=True, color=NEAR_BLACK)
        y += 1.65

    # Constraint comparison
    add_rect(slide, 0.3, y, 12.73, 0.75, fill_rgb=LIGHT_BLUE,
             line_rgb=MID_BLUE, line_width_pt=0.5)
    add_text(slide, "Original constraints:  C1, C1b, C1c, McCormick, C2 SOCP, C3 isolation, C4 control-plane, C5 latency, C6 migration, C7 DRF",
             0.5, y + 0.05, 12.2, 0.32, font_size=12, color=NEAR_BLACK, wrap=True)
    add_text(slide, "Refactored constraints:  C1a capacity, C1b cone, C2 priority link, C3 demand, C4 fairness, C5 node activation",
             0.5, y + 0.37, 12.2, 0.32, font_size=12, bold=True, color=DARK_BLUE, wrap=True)


def slide_complexity(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Complexity Comparison")
    footer(slide, prs, num)

    two_col_table(slide,
        ["Aspect", "Original MISOCP", "Refactored MISOCP"],
        [
            ["Binary variables",
             "O(T·W·N·H + T·N·H + N·H + T)",
             "O(T·N·H + N·H + T)"],
            ["Continuous variables",
             "O(T·W·N·K·H + T·W·H + N·H)",
             "O(T·N·H + N·H + 1)"],
            ["Cone constraints",
             "‖L_r · ξ_{n,r,t}‖₂ ≤ soc_aux  per (n,r,t)",
             "sum_i σ²_{i,h} · y_{i,n,h} ≤ t_{n,h}²  per (n,h)"],
            ["Isolation variables",
             "O(T·W·N·K·H) binary  (McCormick)",
             "None"],
            ["Migration variables",
             "O(T·W·N·H) binary",
             "None"],
            ["External data required",
             "Σ_r,  d_{i,j,r},  N_{i,t},  η_{k,r},  γ_{i,j},  …",
             "u_{i,h}  and  sigma_frac  (one scalar)"],
        ],
        x=0.3, y=1.25, w=12.7, h=4.5,
        col_widths=[3.0, 4.85, 4.85], font_size=13, header_font_size=14)

    callout(slide,
        "The model is simpler to solve but not simpler in class. "
        "MISOCP complexity is retained through the Cantelli cone — "
        "it is just driven by practically obtainable inputs.",
        x=0.3, y=5.9, w=12.7, h=0.82, font_size=14)


# ── Part 2 slides ──────────────────────────────────────────────────────────

def slide_architecture(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "System Architecture", "How the three layers connect")
    footer(slide, prs, num)

    layers = [
        ("PREDICTION LAYER  (teammates)",
         "Time-series forecast  →  u[i,h]  per tenant per planning period",
         LIGHT_BLUE, MID_BLUE),
        ("PLAN-AHEAD MODEL  (MISOCP — Gurobi)",
         "Solved once per planning horizon  →  TenantAccessSchedule (priority hints  y[i,n,h])",
         _rgb(0xE6,0xF4,0xEA), GREEN_FG),
        ("REAL-TIME MODEL  (MILP — OR-Tools / CBC)",
         "Solved every scheduling epoch  →  placement decisions  x[j,n]",
         _rgb(0xFF,0xF3,0xCD), _rgb(0x85,0x65,0x04)),
        ("SIMULATION  (FastAPI + React)",
         "Interactive step-by-step execution with live dashboard",
         LIGHT_GRAY, DARK_BLUE),
    ]

    y = 1.28
    for i, (title, desc, bg, fg) in enumerate(layers):
        add_rect(slide, 0.5, y, 12.33, 1.12, fill_rgb=bg,
                 line_rgb=fg, line_width_pt=1.0)
        add_rect(slide, 0.5, y, 0.08, 1.12, fill_rgb=fg)
        add_text(slide, title, 0.72, y + 0.08, 11.5, 0.38,
                 font_size=15, bold=True, color=fg)
        add_text(slide, desc, 0.72, y + 0.48, 11.5, 0.55,
                 font_size=13, color=NEAR_BLACK, wrap=True)
        y += 1.12
        if i < len(layers) - 1:
            add_text(slide, "▼", 6.4, y, 0.6, 0.28,
                     font_size=16, color=DARK_BLUE, align=PP_ALIGN.CENTER)
            y += 0.28

    add_text(slide,
        "Plug-in point: replace  build_synthetic_data()  with prediction model output. "
        "Interface contract:  dict[(tenant_id, period) → float]",
        0.3, 6.85, 12.7, 0.38, font_size=13, italic=True, color=MID_BLUE, wrap=True)


def slide_plan_ahead_status(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Plan-Ahead Model — Current State")
    footer(slide, prs, num)

    status_badge(slide, "✓  Implemented", 0.3, 1.28)
    status_badge(slide, "✓  Tested", 1.3, 1.28)
    add_text(slide, "plan_ahead_optimizer.py  (Gurobi MISOCP)",
             2.35, 1.3, 10.0, 0.28, font_size=14, italic=True, color=MID_BLUE)

    two_col_table(slide,
        ["Variable", "Description"],
        [
            ["f[i,n,h] >= 0",         "Continuous allocation of node n capacity to tenant i in period h"],
            ["y[i,n,h] in {0,1}",     "Priority indicator — primary output; forms TenantAccessSchedule"],
            ["z[n,h]   in {0,1}",     "Node activation indicator"],
            ["t[n,h]   >= 0",         "Cantelli slack (SOCP mode only)"],
            ["a[i]     in {0,1}",     "Tenant admission indicator"],
            ["sigma    in [0,1]",     "Fairness auxiliary: min demand-satisfaction ratio"],
        ],
        x=0.3, y=1.72, w=12.7, h=3.0,
        col_widths=[3.2, 9.5], font_size=13, header_font_size=14)

    bullet_box(slide,
        ["C1a: capacity with safety buffer",
         "C1b: cone — sum_i σ²_{i,h} · y_{i,n,h} ≤ t_{n,h}²",
         "C2: priority link (f ≤ C[n] · y)",
         "C3: demand satisfaction  (sum_n f ≥ u · a)",
         "C4: fairness (min demand-satisfaction ratio)",
         "C5: node activation  (z ≥ y)"],
        x=0.3, y=4.85, w=6.0, h=2.0,
        bg=LIGHT_GRAY, title="Constraints", font_size=13)

    bullet_box(slide,
        ["Output: TenantAccessSchedule  —  dict[(i,h) → [node_ids]]",
         "use_socp=True  (default): full MISOCP, probabilistic guarantee",
         "use_socp=False: plain MILP, simulation-speed fallback",
         "Sensitivity analysis complete  (plan_ahead_sensitivity.py)"],
        x=6.5, y=4.85, w=6.5, h=2.0,
        bg=LIGHT_BLUE, title="Notes", font_size=13)


def slide_realtime_status(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Real-Time Model — Current State")
    footer(slide, prs, num)

    status_badge(slide, "✓  Implemented", 0.3, 1.28)
    status_badge(slide, "✓  Tested", 1.3, 1.28)
    add_text(slide, "optimizer_google_or.py  (OR-Tools / CBC)",
             2.35, 1.3, 10.0, 0.28, font_size=14, italic=True, color=MID_BLUE)

    add_text(slide, "Objective (maximise):", 0.3, 1.72, 12.7, 0.35,
             font_size=15, bold=True, color=DARK_BLUE)
    add_text(slide,
        "max  sum_{j,n}  omega_{delay,t(j)} · P_j_mem · u_n_mem · sigma_n_consolid "
        "· b_{t(j),n} · x[j,n]",
        0.3, 2.07, 12.7, 0.38, font_size=14, wrap=True)
    add_text(slide, "where  b_{t(j),n} = 2.0  if n in plan-ahead priority set,  else  1.0",
             0.3, 2.45, 12.7, 0.3, font_size=13, italic=True)

    two_col_table(slide,
        ["Constraint", "Description"],
        [
            ["C1", "At most one node per job"],
            ["C2", "Memory capacity ≤ M_n_eff = M_n_cap - U_n - v_bar_n · M_n_cap"],
            ["C4", "CPU fitment pre-filter (x[j,n]=0 if job CPU > node cores)"],
            ["C5 (soft)", "Priority boost: b=2.0 for plan-ahead-endorsed (tenant, node) pairs"],
        ],
        x=0.3, y=2.9, w=12.7, h=2.1,
        col_widths=[1.2, 11.5], font_size=13, header_font_size=14)

    bullet_box(slide,
        ["SLA violation rate v_bar_n: shrinks effective capacity on struggling nodes",
         "Tenant delay weight omega_{delay,t}: boosts priority for tenants with long queue waits",
         "Both updated every scheduling epoch — continuous closed-loop feedback"],
        x=0.3, y=5.15, w=12.7, h=1.55,
        bg=LIGHT_GRAY, title="Active Feedback Loops", font_size=13)


def slide_pipeline(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Pipeline — How the Layers Connect")
    footer(slide, prs, num)

    status_badge(slide, "✓  Implemented", 0.3, 1.28)
    add_text(slide, "Pipeline/interface.py", 1.35, 1.3, 10.0, 0.28,
             font_size=14, italic=True, color=MID_BLUE)

    steps = [
        ("1", "build_synthetic_data()",
         "Generates u[i,h] as placeholders ← PLUG-IN POINT for prediction team output"),
        ("2", "Plan-Ahead MISOCP (Gurobi)",
         "Solves over horizon H → TenantAccessSchedule:  dict[(tenant_id, h) → [node_ids]]"),
        ("3", "schedule_to_leases()",
         "Compresses contiguous same-node periods into TenantLease objects (cleaner period tracking)"),
        ("4", "filter_active_access(schedule, current_period)",
         "Slices TenantAccessSchedule to current planning period → dict[tenant_id → [node_ids]]"),
        ("5", "Real-Time Model (OR-Tools)",
         "Receives active priority hints, applies PRIORITY_BOOST=2.0, places jobs, returns node states"),
    ]

    y = 1.72
    step_h = 0.95
    for step_num, title, desc in steps:
        add_rect(slide, 0.3, y, 0.6, step_h, fill_rgb=DARK_BLUE)
        add_text(slide, step_num, 0.3, y + 0.2, 0.6, 0.55,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, 0.9, y, 12.13, step_h, fill_rgb=LIGHT_GRAY,
                 line_rgb=MED_GRAY, line_width_pt=0.3)
        add_text(slide, title, 1.05, y + 0.07, 11.8, 0.35,
                 font_size=14, bold=True, color=DARK_BLUE)
        add_text(slide, desc, 1.05, y + 0.45, 11.8, 0.4,
                 font_size=13, color=NEAR_BLACK, wrap=True)
        y += step_h + 0.05

    callout(slide,
        "Plug-in contract for prediction team:  dict[(tenant_id: int, period: int) → float]\n"
        "Assign to  P['u']  before calling  build_model(P, env).  No other changes required.",
        x=0.3, y=6.88, w=12.7, h=0.58, font_size=13)


def slide_simulation_status(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Simulation — Current State")
    footer(slide, prs, num)

    status_badge(slide, "✓  Running", 0.3, 1.28)
    add_text(slide, "FastAPI backend  +  React frontend  (Simulation/)",
             1.35, 1.3, 10.0, 0.28, font_size=14, italic=True, color=MID_BLUE)

    bullet_box(slide,
        ["Manages a live job queue with per-tenant arrival rates",
         "Runs the real-time MILP every step (configurable step size)",
         "Fires the plan-ahead MISOCP every plan_ahead_interval steps (default: 50)",
         "Tracks per-node memory usage, SLA violation rates, tenant wait times",
         "REST API: GET /api/state,  POST /api/step,  POST /api/reset,  POST /api/config"],
        x=0.3, y=1.72, w=6.1, h=3.0,
        bg=LIGHT_GRAY, title="What It Does", font_size=13)

    bullet_box(slide,
        ["Node utilization over time (rolling history)",
         "Queue depth per tenant",
         "Placed vs rejected jobs per batch",
         "Plan-ahead priority assignments (tenant → nodes)",
         "SLA violation rate per node"],
        x=6.6, y=1.72, w=6.43, h=3.0,
        bg=LIGHT_BLUE, title="Dashboard Shows (live)", font_size=13)

    bullet_box(slide,
        ["Per-job lifetimes replace old fractional memory release — accurate per-job lifecycle",
         "Actual memory = predicted × (1 + spike_frac) — tests robustness to demand spikes",
         "Simulated clock advances BATCH_DURATION_SEC per step",
         "Node used_mb recomputed from scratch each round (always accurate, self-correcting)"],
        x=0.3, y=4.87, w=12.73, h=1.9,
        bg=LIGHT_GRAY, title="Simulation Mechanics", font_size=13)


def slide_whats_next(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "What’s Next", "Immediate next steps")
    footer(slide, prs, num)

    add_text(slide, "Primary open item: Plug in Prediction Model Output",
             0.3, 1.28, 12.7, 0.42, font_size=20, bold=True, color=DARK_BLUE)

    add_rect(slide, 0.3, 1.78, 12.73, 2.5, fill_rgb=LIGHT_GRAY,
             line_rgb=MED_GRAY, line_width_pt=0.5)
    add_text(slide, "The optimization layer is ready. The interface contract is defined:",
             0.5, 1.85, 12.2, 0.35, font_size=14, wrap=True)
    add_text(slide,
        "# Replace this:\n"
        "P = build_synthetic_data()           # synthetic u[i,h]\n\n"
        "# With this:\n"
        "P['u'] = prediction_model.forecast(tenants, horizon)   # dict[(i,h) → float]",
        0.5, 2.23, 12.2, 1.85, font_size=13, color=_rgb(0x1A,0x3A,0x5C), wrap=True,
        font_name="Courier New")

    add_text(slide,
        "The prediction team delivers:   dict[(tenant_id: int, period: int) → total_resource_usage: float]",
        0.5, 4.35, 12.2, 0.38, font_size=14, bold=True, color=MID_BLUE, wrap=True)
    add_text(slide, "No other changes required in the optimization layer.",
             0.5, 4.73, 12.2, 0.35, font_size=14, color=NEAR_BLACK)

    bullet_box(slide,
        ["Sensitivity analysis already complete — lambda weights, sigma_frac, kappa documented",
         "Simulation tested end-to-end with synthetic usage profiles",
         "Plan-ahead fires every 50 steps; real-time fires every step — verified pipeline timing"],
        x=0.3, y=5.25, w=12.73, h=1.85,
        bg=LIGHT_BLUE, title="Already Verified", font_size=13)


def slide_future_work(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Future Work & Research Directions")
    footer(slide, prs, num)

    # Heuristic box
    add_rect(slide, 0.3, 1.25, 12.73, 3.1, fill_rgb=LIGHT_GRAY,
             line_rgb=MED_GRAY, line_width_pt=0.5)
    add_rect(slide, 0.3, 1.25, 0.08, 3.1, fill_rgb=DARK_BLUE)
    add_text(slide, "Heuristic Acceleration for the Real-Time Model  (exploratory)",
             0.52, 1.3, 12.0, 0.4, font_size=16, bold=True, color=DARK_BLUE)
    add_text(slide,
        "The real-time MILP (OR-Tools / CBC) is exact and fast for small instances "
        "(≤20 jobs, 5 nodes). For larger instances, alternatives include:\n\n"
        "  •  Greedy First-Fit Decreasing (FFD) — O(J log J), well-studied for bin packing. "
        "Can be adapted to respect the plan-ahead priority boost. Likely sufficient "
        "for most practical cluster sizes.\n\n"
        "  •  LP relaxation + rounding (OR-Tools GLOP mode) — faster, slight quality loss.\n\n"
        "  •  Learning-based heuristics — policy trained on historical placements; "
        "out of scope for capstone timeline.",
        0.52, 1.73, 12.1, 2.5, font_size=13, wrap=True)

    add_rect(slide, 0.3, 4.5, 12.73, 0.55, fill_rgb=RED_BG,
             line_rgb=RED_FG, line_width_pt=0.8)
    add_text(slide,
        "Capstone scope note: heuristic acceleration is exploratory. "
        "A working exact MILP is already in place. "
        "If time permits we will benchmark FFD vs CBC; otherwise this is future work.",
        0.5, 4.56, 12.2, 0.42, font_size=13, color=RED_FG, wrap=True)

    bullet_box(slide,
        ["Replace synthetic u_{i,h} with live Google cluster trace replay",
         "Tune Cantelli epsilon per-tenant (risk-aware SLA contracts)",
         "Multi-resource extension: CPU + memory jointly in the plan-ahead cone",
         "Benchmark plan-ahead solve time vs horizon length and tenant count"],
        x=0.3, y=5.22, w=12.73, h=1.88,
        bg=LIGHT_GRAY, title="Other Future Directions", font_size=13)


def slide_summary(prs, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Summary")
    footer(slide, prs, num)

    two_col_table(slide,
        ["Layer", "Status", "Next Step"],
        [
            ["Plan-Ahead (MISOCP)",  "✅  Implemented & tested",    "Swap synthetic u[i,h] for prediction output"],
            ["Real-Time (MILP)",     "✅  Implemented & tested",    "—"],
            ["Pipeline",             "✅  End-to-end working",      "—"],
            ["Simulation",           "✅  Running with dashboard",  "—"],
            ["Prediction integration","⏳  Waiting on teammates",   "Define handoff format, plug in"],
            ["Heuristic RT solver",  "\U0001f52c  Exploratory",         "Future work beyond capstone"],
        ],
        x=0.3, y=1.25, w=12.7, h=4.0,
        col_widths=[3.2, 3.5, 6.0], font_size=14, header_font_size=15)

    callout(slide,
        "The optimization layer is complete and functional as a standalone system. "
        "The primary open integration point is the prediction model output  u[i,h].",
        x=0.3, y=5.45, w=12.7, h=0.82, font_size=15)

    add_text(slide,
        "Plan-ahead refactoring preserves MISOCP complexity while replacing unknowable "
        "inputs (per-workload covariance Σ_r) with forecastable ones (aggregate tenant demand u_{i,h}).",
        0.3, 6.42, 12.7, 0.65, font_size=14, italic=True, color=MID_BLUE, wrap=True)


def slide_questions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
    add_rect(slide, 0, 3.4, 13.33, 0.06, fill_rgb=ACCENT)
    add_text(slide, "Questions?", 0, 2.3, 13.33, 0.9,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Alrick Grandison  |  Capstone Project, Spring 2026",
             0, 3.7, 13.33, 0.5, font_size=20, color=MED_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide,
        "Cluster_Optimization_Models / Docs / plan_ahead_refactoring.pptx",
        0, 4.3, 13.33, 0.4, font_size=14,
        color=_rgb(0x88, 0xA8, 0xC8), align=PP_ALIGN.CENTER)


# ── Assemble ───────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # blank layout for all slides
    blank = prs.slide_layouts[6]

    slide_title(prs)                                        # 1
    slide_agenda(prs, 2)                                    # 2

    section_divider(prs, 1,
        "Plan-Ahead Refactoring",
        "Justification for design changes to the planning model")     # 3

    slide_original_model(prs, 4)                            # 4
    slide_problems(prs, 5)                                  # 5
    slide_change1(prs, 6)                                   # 6
    slide_change2(prs, 7)                                   # 7
    slide_change3(prs, 8)                                   # 8
    slide_misocp_preserved(prs, 9)                          # 9
    slide_removed(prs, 10)                                  # 10
    slide_complexity(prs, 11)                               # 11

    section_divider(prs, 2,
        "Optimization Layer Status",
        "Current state of all three model layers")                    # 12

    slide_architecture(prs, 13)                             # 13
    slide_plan_ahead_status(prs, 14)                        # 14
    slide_realtime_status(prs, 15)                          # 15
    slide_pipeline(prs, 16)                                 # 16
    slide_simulation_status(prs, 17)                        # 17
    slide_whats_next(prs, 18)                               # 18
    slide_future_work(prs, 19)                              # 19
    slide_summary(prs, 20)                                  # 20
    slide_questions(prs)                                    # 21

    out = "Cluster_Optimization_Models/Docs/plan_ahead_refactoring.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
