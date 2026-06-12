"""
fill_optimization.py
────────────────────
Fills ALRICK's "Optimization Models" column in the group master poster with two
UML-style process flowcharts (one per model), showing how each model works:

  • Real-Time Optimizer (MILP)        — per-interval job → machine placement, with
                                        a "Placed?" decision branch.
  • Plan-Ahead Optimizer (MISOCP)     — forecast → tag tenants → solve → schedule,
                                        feeding the real-time layer.

Reads the pristine template (PosterSamples/POSTER.pptx) and writes the filled
copy to PosterComponents/POSTER.pptx. Nothing in the prediction team's columns
is touched. The shared "Limitations & Future Research" card is left untouched.

Flow content is grounded in the source of truth:
  Realtime/realtime_optimizer.py, Realtime/cluster_manager.py,
  PlanAhead/plan_ahead_optimizer.py

Run:  python fill_optimization.py
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

HERE = Path(__file__).resolve().parent
SRC  = HERE.parent / "PosterSamples" / "POSTER.pptx"
OUT  = HERE / "POSTER.pptx"

# ── palette (harmonised with the template) ──────────────────────────────
NAVY  = "0C4480"
LIGHT = "F3F4F0"
DARK  = "1A3A6B"
BLUE  = "2C6EBD"
TEAL  = "14B8A6"
AMBER = "F2A900"
MUTE  = "9DB3CC"
PROC_FILL = "E9EFF7"; PROC_LINE = "C9D6E8"
YES_FILL  = "E4F5F0"; YES_LINE  = TEAL;  YES_TXT = "0B5A50"
NO_FILL   = "FBEEDA"; NO_LINE   = AMBER; NO_TXT  = "6B4E12"
DIA_TXT   = "3A2A00"

def C(s): return RGBColor.from_string(s)

prs = Presentation(str(SRC))
slide = prs.slides[0]
shapes = slide.shapes
by_id = {sh.shape_id: sh for sh in shapes}

# ── primitives ──────────────────────────────────────────────────────────
def box(x, y, w, h, fill=None, line=None, lw=1.0,
        shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def words(sp, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          font="Open Sans"):
    """lines = list of (text, size, bold, color[, italic])"""
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        t, sz, bold, col = ln[0], ln[1], ln[2], ln[3]
        ital = ln[4] if len(ln) > 4 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = 1.0; p.space_after = Pt(1)
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
        r.font.name = font; r.font.color.rgb = C(col)
    return sp

def label(x, y, w, h, text, size=11, color=LIGHT, bold=True,
          align=PP_ALIGN.CENTER, font="Open Sans"):
    tb = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(1); tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = C(color)
    return tb

# ── flowchart node helpers ──────────────────────────────────────────────
def proc(x, y, w, h, lines, fill=PROC_FILL, line=PROC_LINE, lw=1.25,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = box(x, y, w, h, fill=fill, line=line, lw=lw, shape=shape, radius=radius)
    words(sp, lines)
    return sp

def pill(x, y, w, h, lines, fill, txt):
    sp = box(x, y, w, h, fill=fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.5)
    words(sp, [(t, s, b, txt) + tuple(rest) for (t, s, b, *rest) in
               [(l[0], l[1], l[2], *l[3:]) for l in lines]])
    return sp

def solvebox(x, y, w, h, lines):
    sp = box(x, y, w, h, fill=LIGHT, line=AMBER, lw=2.5,
             shape=MSO_SHAPE.RECTANGLE)
    words(sp, lines)
    return sp

def diamond(cx, y, w, h, text):
    sp = box(cx - w/2, y, w, h, fill=AMBER, line=None, shape=MSO_SHAPE.DIAMOND)
    words(sp, [(text, 14.5, True, DIA_TXT)])
    return sp

def darrow(cx, y, h=0.24, w=0.42, color=MUTE):
    box(cx - w/2, y, w, h, fill=color, line=None, shape=MSO_SHAPE.DOWN_ARROW)

def rarrow(x, cy, w=0.55, h=0.34, color=MUTE):
    box(x, cy - h/2, w, h, fill=color, line=None, shape=MSO_SHAPE.RIGHT_ARROW)

# ── geometry of the optimization column ─────────────────────────────────
CX  = 29.80                 # column centre x
BW  = 9.6                   # standard node width
BX  = CX - BW/2             # 25.0

# ════════════════════════════════════════════════════════════════════════
#  Sub-title 1 : Real-Time Optimizer
# ════════════════════════════════════════════════════════════════════════
rt = by_id[70]
rt.left, rt.top, rt.width, rt.height = (Inches(24.15), Inches(20.48),
                                        Inches(11.3), Inches(0.95))
rt.text_frame.clear()
words(rt, [("Real-Time Optimizer", 27, True, LIGHT),
           ("Mixed-Integer Linear Program (MILP)", 14, False, MUTE, True)],
      font="Source Sans Pro")

# ── bubble colours ──────────────────────────────────────────────────────
DELIV = "0E9384"   # teal   (intake)
FILT  = "2C6EBD"   # blue   (prepare)
OPTB  = "F2A900"   # amber  (the optimizer — dark text)
OPTTX = "3A2A00"
DEC   = "6C5CE7"   # purple (decision)
NOB   = "E0663B"   # orange (no)
YESB  = "1E9E6A"   # green  (yes / output)
ENDB  = "1F4E96"   # deep blue (handoff)

# ── Real-Time flowchart (bubbles) ───────────────────────────────────────
pill(BX, 21.55, BW, 0.88,
     [("Cluster Manager delivers jobs to the pending queue", 14.5, True)],
     DELIV, LIGHT)
darrow(CX, 22.45)
pill(BX, 22.69, BW, 0.88,
     [("Filter the queue & machines for each tenant group", 14.5, True)],
     FILT, LIGHT)
darrow(CX, 23.59)
pill(BX, 23.83, BW, 1.02,
     [("Real-Time MILP", 17, True),
      ("assigns each job to a node", 13, False)], OPTB, OPTTX)
darrow(CX, 24.91)
# decision bubble
DCX = 28.05
pill(26.35, 25.21, 3.4, 0.98, [("Assigned?", 16, True)], DEC, LIGHT)
# No → right
rarrow(29.85, 25.70, w=0.6)
label(29.05, 25.30, 1.3, 0.3, "no", size=11, color=MUTE)
pill(30.55, 25.18, 4.85, 1.04,
     [("No — stays in the queue;", 12.5, True),
      ("wait time grows, retried next interval", 11.5, False)], NOB, LIGHT)
# Yes → down
darrow(DCX, 26.19)
label(DCX + 0.30, 26.22, 1.0, 0.3, "yes", size=11, color=MUTE,
      align=PP_ALIGN.LEFT)
pill(DCX - 3.5, 26.45, 7.0, 0.95,
     [("Yes — Cluster Manager places the job on the node", 13, True)],
     YESB, LIGHT)

# ════════════════════════════════════════════════════════════════════════
#  Sub-title 2 : Plan-Ahead Optimizer
# ════════════════════════════════════════════════════════════════════════
pa = by_id[71]
pa.left, pa.top, pa.width, pa.height = (Inches(24.15), Inches(27.92),
                                        Inches(11.3), Inches(0.95))
pa.text_frame.clear()
words(pa, [("Plan-Ahead Optimizer", 27, True, LIGHT),
           ("Chance-Constrained MISOCP", 14, False, MUTE, True)],
      font="Source Sans Pro")

# ── Plan-Ahead flowchart (bubbles) ──────────────────────────────────────
pill(BX, 29.00, BW, 0.95,
     [("Forecast each tenant's demand over the horizon", 14.5, True),
      ("mean + variance per period", 12, False)], DELIV, LIGHT)
darrow(CX, 30.00)
pill(BX, 30.24, BW, 0.88,
     [("Tag tenants: exclusive vs shared", 15, True)], FILT, LIGHT)
darrow(CX, 31.14)
pill(BX, 31.38, BW, 1.45,
     [("Plan-Ahead MISOCP", 17, True),
      ("reserves machines for each tenant per period", 12.5, False),
      ("Cantelli cone keeps  P(overflow) ≤ ε", 12.5, False)], OPTB, OPTTX)
darrow(CX, 32.86)
pill(BX, 33.10, BW, 0.88,
     [("Per-period schedule:  tenant groups → machines", 14, True)], YESB, LIGHT)
darrow(CX, 34.00)
pill(BX, 34.24, BW, 0.82,
     [("Feeds the tenant groups to the Real-Time Optimizer", 13.5, True)], ENDB, LIGHT)

# ── save ────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print("Wrote:", OUT)
print("Total shapes now:", len(shapes._spTree))
